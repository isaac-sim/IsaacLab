"""Build the compact (14-slide) sensor simulation release deck.

This is the focused variant matching the structure of the 12-slide PDF
(sensor_simulation_release_deck (3).pptx), updated with all XCATH
requirements slides and vessel mesh collision features.

Usage:
    python build_compact_deck.py
"""

from __future__ import annotations

import sys
from pathlib import Path

# Import everything from the full-deck builder so we reuse helpers/builders.
sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_release_deck import (  # noqa: E402
    # theme
    NV_GREEN,
    NV_GREEN_DARK,
    DARK_BG,
    LIGHT_BG,
    WHITE,
    BLACK,
    GREY,
    LIGHT_GREY,
    DONE_GREEN,
    PARTIAL_ORANGE,
    PLANNED_RED,
    SLIDE_W,
    SLIDE_H,
    # helpers
    add_rect,
    add_text,
    add_header_bar,
    add_bullets,
    add_table,
    add_footer,
    make_content_slide,
    make_section_divider,
    # slide builders reused verbatim
    make_title_slide,
    slide_xcath_requirements,
    slide_xcath_progress,
    slide_modality_status,
    slide_xray_perf,
    slide_rl_pipeline,
    slide_sprint1_completed,
    slide_closed_gaps,
    slide_xcath_rod_solver,
    slide_next_xray,
    slide_summary,
    slide_questions,
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Compact-only slide: architecture overview (matches slide 2 of the PDF)
# ---------------------------------------------------------------------------


def slide_architecture_compact(prs):
    """High-level system architecture — compact single-slide overview."""
    slide = make_content_slide(
        prs,
        "Architecture Overview",
        subtitle="CT ingestion → physics simulation → X-ray render → RL training loop",
    )

    # Row 1: pipeline stages
    stages = [
        ("CT Ingestion", "DICOM/NIfTI → HU→μ → vessel seg → centerline → arrival map → SDF → USD"),
        ("Physics Simulation", "XPBD Cosserat rod (3 backends: PyTorch prod / self-contained Warp / Newton bridge)  ·  512 envs  ·  CUDA graphs"),
        ("X-Ray Rendering", "Slang DiffDRR fused volume + catheter  ·  Beer-Lambert compositing  ·  DSA pipeline  ·  9 C-arm presets"),
        ("RL Training", "RSL-RL PPO  ·  state observations (Sprint 1)  ·  pixel/fluoro obs (Sprint 2)  ·  domain randomization"),
    ]

    box_w = (SLIDE_W - Inches(1.2)) / 4
    box_h = Inches(1.5)
    y_top = Inches(1.85)

    for i, (title, desc) in enumerate(stages):
        x = Inches(0.5) + i * (box_w + Inches(0.07))
        add_rect(slide, x, y_top, box_w, box_h, NV_GREEN_DARK)
        add_text(slide, title, x + Inches(0.08), y_top + Inches(0.05), box_w - Inches(0.16), Inches(0.4),
                 size=13, bold=True, color=WHITE)
        add_text(slide, desc, x + Inches(0.08), y_top + Inches(0.45), box_w - Inches(0.16), Inches(1.0),
                 size=10, color=LIGHT_GREY)
        # arrow between boxes
        if i < 3:
            ax = x + box_w + Inches(0.01)
            add_text(slide, "→", ax, y_top + Inches(0.55), Inches(0.06), Inches(0.4),
                     size=18, bold=True, color=NV_GREEN)

    # Row 2: key data flows
    add_text(slide, "Key Data Flows", Inches(0.5), Inches(3.55), SLIDE_W - Inches(1.0), Inches(0.35),
             size=14, bold=True, color=NV_GREEN_DARK)
    flows = [
        "CTA volume (μ array, 512³) ─────────────────────────────────────────────→ Slang ray-marcher",
        "Vessel mesh + SDF ──────────────────────────────────────────────────────→ XCathRodSolver BVH collision",
        "Rod positions (Warp tensor) ─────────────────────────────────────────────→ Slang catheter compositing",
        "Fluoroscopy frame (H×W float) ──────────────────────────────────────────→ RL observation dict (Sprint 2)",
        "Reward signal ───────────────────────────────────────────────────────────→ RSL-RL PPO update",
    ]
    add_bullets(slide, [(f, False) for f in flows],
                Inches(0.5), Inches(4.0), SLIDE_W - Inches(1.0), Inches(2.2),
                size=12, line_spacing=1.3)

    # Status bar
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.35), LIGHT_BG)
    add_text(slide,
             "Sprint 1 complete: CT ingestion → physics → X-ray render → state RL.  "
             "Sprint 2: fluoroscopy pixel obs + GPU detector physics + multi-env vessel collision.",
             Inches(0.4), Inches(7.12), SLIDE_W - Inches(0.8), Inches(0.3),
             size=10, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 14-slide compact sequence
    builders = [
        # ── Context ──────────────────────────────────────────────────────
        ("title",         make_title_slide),
        ("arch",          slide_architecture_compact),
        ("at_a_glance",   slide_modality_status),          # full component table incl. vessel collision row
        # ── XCATH partner requirements (2 new slides) ────────────────────
        ("xcath_req",     slide_xcath_requirements),
        ("xcath_prog",    slide_xcath_progress),
        # ── Current status ───────────────────────────────────────────────
        ("part1",         lambda p: make_section_divider(
            p, "PART 1", "Current Status",
            "What is built, integrated, and measured today."
        )),
        ("xray_perf",     slide_xray_perf),
        ("rl_pipeline",   slide_rl_pipeline),
        # ── This release ─────────────────────────────────────────────────
        ("part2",         lambda p: make_section_divider(
            p, "PART 2", "This Release",
            "Sprint 1 — features shipping now."
        )),
        ("deliverables",  slide_sprint1_completed),        # incl. XCathRodSolver rows
        ("closed_gaps",   slide_closed_gaps),              # incl. vessel mesh collision rows
        ("xcath_solver",  slide_xcath_rod_solver),         # XCathRodSolver detail
        # ── Next release ─────────────────────────────────────────────────
        ("part3",         lambda p: make_section_divider(
            p, "PART 3", "Next Release",
            "Sprint 2 — Texture2DArray, multi-env vessel collision, image observations."
        )),
        ("next_xray",     slide_next_xray),
        # ── Summary ──────────────────────────────────────────────────────
        ("summary",       slide_summary),
        ("qna",           slide_questions),
    ]

    for _, fn in builders:
        fn(prs)

    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1 or idx == total:
            continue
        add_footer(slide, idx, total)

    out = Path(__file__).resolve().parent / "sensor_simulation_release_deck_compact.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    build()
