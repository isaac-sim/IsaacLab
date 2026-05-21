"""Build the Sensor Simulation release-status PowerPoint deck.

Generates `sensor_simulation_release_deck.pptx` in the same directory.
Usage:
    python build_release_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

NV_GREEN = RGBColor(0x76, 0xB9, 0x00)
NV_GREEN_DARK = RGBColor(0x4F, 0x7A, 0x00)
DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY = RGBColor(0xE0, 0xE0, 0xE0)

DONE_GREEN = RGBColor(0x2E, 0x7D, 0x32)
PARTIAL_ORANGE = RGBColor(0xEF, 0x6C, 0x00)
PLANNED_RED = RGBColor(0xB7, 0x1C, 0x1C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Calibri",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.55), NV_GREEN)
    add_text(
        slide,
        "X-Ray–Guided Robotic Catheter Intervention — Isaac for Healthcare",
        Inches(0.4),
        Inches(0.08),
        Inches(9),
        Inches(0.4),
        size=14,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        "NVIDIA Healthcare — Holoscan Team",
        Inches(8.5),
        Inches(0.08),
        Inches(4.5),
        Inches(0.4),
        size=12,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        title,
        Inches(0.5),
        Inches(0.75),
        SLIDE_W - Inches(1.0),
        Inches(0.6),
        size=30,
        bold=True,
        color=NV_GREEN_DARK,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.5),
            Inches(1.30),
            SLIDE_W - Inches(1.0),
            Inches(0.4),
            size=16,
            color=GREY,
        )
    add_rect(slide, Inches(0.5), Inches(1.32 if not subtitle else 1.72), Inches(1.0), Inches(0.04), NV_GREEN)


def add_footer(slide, page_num, total):
    add_text(
        slide,
        f"NVIDIA Healthcare — Holoscan Team   ·   {page_num} / {total}",
        Inches(0.4),
        SLIDE_H - Inches(0.35),
        SLIDE_W - Inches(0.8),
        Inches(0.3),
        size=10,
        color=GREY,
        align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, items, x, y, w, h, *, size=18, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"•  {text}"
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = BLACK
    return tb


def style_table(table, *, header_fill=NV_GREEN, header_color=WHITE, body_size=12, header_size=13):
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    for c in range(n_cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(header_size)
                r.font.color.rgb = header_color
                r.font.name = "Calibri"
    for r_idx in range(1, n_rows):
        zebra = LIGHT_BG if r_idx % 2 == 1 else WHITE
        for c in range(n_cols):
            cell = table.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = zebra
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(body_size)
                    r.font.name = "Calibri"
                    if r.font.color.rgb is None:
                        r.font.color.rgb = BLACK
    for row in table.rows:
        row.height = Inches(0.0)
    for cell in [c for row in table.rows for c in row.cells]:
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)


STATUS_TOKENS = {
    "[OK]": ("✓", DONE_GREEN, True),
    "[PARTIAL]": ("◐", PARTIAL_ORANGE, True),
    "[PLAN]": ("✗", PLANNED_RED, True),
}


def set_cell_text(cell, text, *, bold=False, color=None, align=PP_ALIGN.LEFT, size=12):
    """Write `text` into a cell; supports inline status tokens and **bold** spans."""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align

    tokens = []
    remaining = text
    while remaining:
        # Find earliest special token
        idx_status = min(
            (remaining.find(k) for k in STATUS_TOKENS if remaining.find(k) != -1),
            default=-1,
        )
        idx_bold = remaining.find("**")
        candidates = [i for i in (idx_status, idx_bold) if i != -1]
        if not candidates:
            tokens.append(("plain", remaining))
            break
        first = min(candidates)
        if first > 0:
            tokens.append(("plain", remaining[:first]))
            remaining = remaining[first:]
            continue
        if remaining.startswith("**"):
            end = remaining.find("**", 2)
            if end == -1:
                tokens.append(("plain", remaining))
                break
            tokens.append(("bold", remaining[2:end]))
            remaining = remaining[end + 2:]
            continue
        for k, (sym, col, b) in STATUS_TOKENS.items():
            if remaining.startswith(k):
                tokens.append(("status", (sym, col)))
                remaining = remaining[len(k):]
                break

    for kind, payload in tokens:
        run = p.add_run()
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        if kind == "status":
            sym, col = payload
            run.text = f"{sym} "
            run.font.bold = True
            run.font.color.rgb = col
        elif kind == "bold":
            run.text = payload
            run.font.bold = True
            run.font.color.rgb = color or BLACK
        else:
            run.text = payload
            run.font.bold = bold
            run.font.color.rgb = color or BLACK


def add_table(slide, x, y, w, h, headers, rows, *, col_widths=None, body_size=12, header_size=13):
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, frac in enumerate(col_widths):
            table.columns[i].width = int(w * frac / total)
    for c, hdr in enumerate(headers):
        set_cell_text(table.cell(0, c), hdr, bold=True, color=WHITE, size=header_size)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            set_cell_text(table.cell(r_idx, c_idx), str(val), size=body_size)
    style_table(table, body_size=body_size, header_size=header_size)
    return table


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    add_rect(slide, 0, Inches(2.4), SLIDE_W, Inches(0.06), NV_GREEN)
    add_rect(slide, 0, Inches(5.0), SLIDE_W, Inches(0.06), NV_GREEN)

    add_text(
        slide,
        "X-Ray–Guided Robotic",
        Inches(0.6),
        Inches(2.55),
        SLIDE_W - Inches(1.2),
        Inches(0.9),
        size=48,
        bold=True,
        color=NV_GREEN,
    )
    add_text(
        slide,
        "Catheter Intervention",
        Inches(0.6),
        Inches(3.20),
        SLIDE_W - Inches(1.2),
        Inches(0.9),
        size=48,
        bold=True,
        color=NV_GREEN,
    )
    add_text(
        slide,
        "Simulation Workflow — Release Status & Roadmap",
        Inches(0.6),
        Inches(4.15),
        SLIDE_W - Inches(1.2),
        Inches(0.6),
        size=24,
        color=WHITE,
    )
    add_text(
        slide,
        "Isaac for Healthcare   ·   NVIDIA Healthcare — Holoscan Team",
        Inches(0.6),
        Inches(5.3),
        SLIDE_W - Inches(1.2),
        Inches(0.5),
        size=16,
        color=LIGHT_GREY,
    )
    return slide


def make_section_divider(prs, label, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    add_rect(slide, Inches(0.5), Inches(2.6), Inches(0.15), Inches(2.5), NV_GREEN)
    add_text(
        slide,
        label,
        Inches(0.85),
        Inches(2.55),
        SLIDE_W - Inches(1.2),
        Inches(0.6),
        size=20,
        bold=True,
        color=NV_GREEN,
    )
    add_text(
        slide,
        title,
        Inches(0.85),
        Inches(3.1),
        SLIDE_W - Inches(1.2),
        Inches(1.2),
        size=44,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.85),
            Inches(4.5),
            SLIDE_W - Inches(1.2),
            Inches(0.6),
            size=18,
            color=LIGHT_GREY,
        )
    return slide


def make_content_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header_bar(slide, title, subtitle)
    return slide


# ---------------------------------------------------------------------------
# Slide content
# ---------------------------------------------------------------------------


def slide_xcath_requirements(prs):
    """XCATH Robotics — what they're building and why they need the simulator."""
    slide = make_content_slide(
        prs,
        "XCATH Robotics — Partner Requirements",
        subtitle="Real-time neurointerventional monitoring system · Jung-eun Park, AI/ML Engineer · March / April 2026",
    )

    # Left column: What XCATH is building
    add_text(
        slide, "What XCATH Is Building",
        Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.4),
        size=14, bold=True,
    )
    left_bullets = [
        "Real-time monitoring system for neurointerventional surgery",
        "  → 3D vessel overlay on live fluoroscopy during catheter procedures",
        "Required AI components:",
        "  1. Vessel segmentation on intra-op fluoroscopy",
        "  2. 2D/3D Registration: pre-op CTA ↔ intra-op fluoroscopy",
        "The Bottleneck:",
        "  Paired training data (CTA + fluoroscopy + GT pose) does NOT",
        "  exist at scale — different modalities, acquisition times, positions",
        "Why NVIDIA's Simulator:",
        "  Input: CTA volume + vessel mask + pose parameters",
        "  Output: Synthetic fluoroscopy / DSA image + exact GT 3D→2D pose",
        "  → Unlimited paired data for registration & segmentation training",
    ]
    add_bullets(
        slide, [(b, False) for b in left_bullets],
        Inches(0.5), Inches(2.45), Inches(6.3), Inches(4.2),
        size=11, line_spacing=1.25,
    )

    # Right column: 5 simulation methods + next steps
    add_text(
        slide, "Simulation Methods Validated (March 2026)",
        Inches(7.0), Inches(1.95), Inches(5.9), Inches(0.4),
        size=14, bold=True,
    )
    right_bullets = [
        "[OK] Method 1: DRR — base fluoroscopy from CTA volume",
        "[OK] Method 2: Vessel Boost — μ × 8 on vessel-masked voxels",
        "[OK] Method 3: DSA — mask subtraction + scatter/misregistration",
        "[OK] Method 5: Bolus Centerline — VMTK + Dijkstra + gamma-variate",
        "",
        "Next steps requested from NVIDIA:",
        "  • Synthetic dataset generation at scale",
        "    (CTA + synth DSA + GT pose → train registration)",
        "  • 2D/3D registration model (PoseNet / DiffDRR architecture)",
        "  • Realism evaluation: FID, SSIM, vessel visibility",
        "  • 50-patient internal CTA+DSA paired dataset available",
        "    (cannot be shared; used for validation only)",
    ]
    add_bullets(
        slide, [(b, False) for b in right_bullets],
        Inches(7.0), Inches(2.45), Inches(5.9), Inches(4.2),
        size=11, line_spacing=1.25,
    )

    add_text(
        slide,
        "Source: XCATH Robotics × NVIDIA Technical Collaboration — March 2026 presentation.",
        Inches(0.5), Inches(6.9), SLIDE_W - Inches(1.0), Inches(0.4),
        size=10, color=GREY,
    )
    return slide


def slide_xcath_progress(prs):
    """XCATH April 2026 progress: validated metrics and remaining gaps."""
    slide = make_content_slide(
        prs,
        "XCATH — April 2026 Progress & Validated Metrics",
        subtitle="Weekly Progress 2026-04 W2–W3 · FluoroSim vs DeepDRR · Selective injection · Dispersion correction",
    )

    headers = ["Area", "March 2026", "April 2026 W1–W3", "Status"]
    rows = [
        ["DSA rendering", "Contrast-only (bone visible)", "Mask subtraction + selective injection (F-8)", "[OK] DONE"],
        ["Blood flow model", "v=150·(r/r_ref)^0.5, allometric", "v=350·(r/2.0)², Hagen-Poiseuille (power=2)", "[OK] DONE"],
        ["Bolus delay match vs real DSA", "5.6 s vs 2.13 s (2.6× off)", "1.90 s vs 2.13 s (89% match)", "[OK] DONE"],
        ["Bolus dispersion", "None — all voxels same shape", "β_eff = β₀ + k·T + α/r  (2.55× asymmetry)", "[OK] DONE"],
        ["Selective injection", "4 roots simultaneous", "Feeding Trunk Excision F-8; 43.1% territory", "[OK] DONE"],
        ["FluoroSim vs DeepDRR", "Qualitative only", "SSIM=0.964, MAE=9.3%, Pearson=0.981", "[OK] DONE"],
        ["Code quality", "Monolithic scripts", "Modular + 125 tests + adapter.py + Dockerfile", "[OK] DONE"],
        ["Multi-patient validation", "Patient #1 only", "2–3 more patients needed", "[PEND] Blocked on data"],
    ]
    add_table(
        slide,
        Inches(0.4), Inches(2.05),
        SLIDE_W - Inches(0.8), Inches(3.8),
        headers, rows,
        col_widths=[20, 26, 36, 18],
        body_size=10, header_size=12,
    )

    add_text(
        slide, "Remaining Gaps Identified by XCATH",
        Inches(0.5), Inches(5.95), SLIDE_W - Inches(1.0), Inches(0.35),
        size=13, bold=True,
    )
    gap_bullets = [
        "Beam hardening (monoenergetic): noise CV sim 0.47 vs real 1.46 (3.1×) — requires polyenergetic rendering",
        "Iodine K-edge (33.2 keV): 70 keV misses 2.5× μ jump → DSA vessel signal weaker than real",
        "Focal-spot blur: point-source renderer → vessels appear crisper than clinical (mitigation: multi-source render)",
        "Finer distal cortical branches: limited by CTA annotation density — not a simulator gap",
    ]
    add_bullets(
        slide, [(b, False) for b in gap_bullets],
        Inches(0.5), Inches(6.35), SLIDE_W - Inches(1.0), Inches(0.95),
        size=10, line_spacing=1.2,
    )
    add_text(
        slide,
        "Key finding: monoenergetic rendering (70 keV) is practically sufficient for cerebral DSA "
        "(beam hardening cancels in mask subtraction; SSIM=0.964 vs DeepDRR polyenergetic). "
        "Source: XCATH Robotics × NVIDIA — April 2026 W2–W3.",
        Inches(0.5), Inches(7.05), SLIDE_W - Inches(1.0), Inches(0.35),
        size=9, color=GREY,
    )
    return slide


def slide_agenda(prs):
    slide = make_content_slide(prs, "Agenda")
    items = [
        ("Executive Snapshot — where the X-Ray catheter workflow is today", True),
        ("Current Status — what's running, with numbers", False),
        ("This Release — features shipping now (Sprint 1)", False),
        ("Next Release — Sprint 2 priorities", False),
        ("Following Releases — Sprint 3 & agentic workflow phases", False),
        ("Workflow Enhancements — adjacent capabilities on the roadmap", False),
    ]
    add_bullets(slide, items, Inches(0.8), Inches(2.0), SLIDE_W - Inches(1.6), Inches(5.0), size=22, line_spacing=1.5)
    return slide


def slide_executive_snapshot(prs):
    slide = make_content_slide(prs, "Executive Snapshot")
    items = [
        ("X-Ray fluoroscopy simulation pipeline validated end-to-end: physics + Beer-Lambert compositing + Slang GPU fused render", False),
        ("Fused volume + catheter Beer-Lambert compositing: ~25 FPS single-env; batched multi-env Slang renderer NOW IMPLEMENTED", False),
        ("All three XPBD rod backends now multi-env: production (PyTorch) + self-contained (Warp) + Newton bridge wrapper", False),
        ("Self-contained XPBD: batched block-Thomas + GPU root control + CUDA-graph capture + floor restitution + 4 solver backends", False),
        ("Full detector physics chain: Poisson + scatter + PSF + gamma + misregistration jitter (CPU path; GPU port = Sprint 2)", False),
        ("Full DSA pipeline (4-step) + vessel boost + gamma-variate bolus dynamics + VMTK/Dijkstra arrival map", False),
        ("Volumetric instrument injection via Warp atomic-max compositing (~2 ms / 64 nodes / 512³)", False),
        ("9 vendor C-arm presets (GE, Siemens, Philips, Ziehm) + state-based RL (PPO, 512 envs)", False),
        ("E2E validated: catheter traverses 35 mm across cranial anatomy in Slang fused render with correct Beer-Lambert polarity", False),
        ("Vessel mesh collision: XCathRodSolver ported from Newton xcath branch — SDF BVH + AABB/edge kernels + track-guided insertion", False),
        ("Critical remaining gap: image-based RL observations + GPU-side detector physics + beam hardening", False),
    ]
    add_bullets(
        slide,
        [(t.replace("**", ""), b) for t, b in items],
        Inches(0.6),
        Inches(1.7),
        SLIDE_W - Inches(1.2),
        Inches(5.5),
        size=17,
        line_spacing=1.35,
    )
    return slide


def slide_modality_status(prs):
    slide = make_content_slide(prs, "X-Ray Fluoroscopy Pipeline — At a Glance")
    headers = ["Component", "Status", "Implementation", "Multi-Env"]
    rows = [
        ["Catheter physics solver", "[OK] Implemented", "Production XPBD Cosserat rod (Warp + Newton bridge)", "[OK] Yes (all 3 backends)"],
        ["DRR volume rendering", "[OK] Implemented", "Slang DiffDRR (differentiable autodiff)", "[OK] Batched (renderDRR_forward_batched)"],
        ["Catheter Beer-Lambert compositing", "[OK] Implemented", "Fused GPU ray-march + CPU detector physics path", "[OK] Batched (StructuredBuffer slice)"],
        ["Volumetric instrument injection", "[OK] Implemented", "Warp atomic-max kernels, ~2 ms / 64 nodes / 512³", "Sprint 2"],
        ["DSA pipeline (4-step)", "[OK] Implemented", "DSAPipeline: contrast/mask DRR + scatter + jitter + boost", "Sprint 2"],
        ["Bolus dynamics (temporal DSA)", "[OK] Implemented", "VMTK centerline + Dijkstra arrival + gamma-variate C(t)", "Sprint 2"],
        ["C-arm geometry + 9 vendor presets", "[OK] Implemented", "GE, Siemens, Philips, Ziehm classmethod factories", "N/A"],
        ["Detector physics chain", "[OK] Implemented", "Poisson + scatter + PSF + gamma + misregistration jitter", "N/A"],
        ["RL training pipeline", "[OK] Implemented", "PPO via RSL-RL, state observations, 512 envs", "[OK] Yes"],
        ["Vessel mesh collision (catheter-in-vessel)", "[OK] Implemented", "XCathRodSolver: SDF BVH + AABB/edge + track guidance (xcath_rod_solver.py)", "Single-env (multi-env = Sprint 2)"],
        ["Image-based RL observations", "[PLAN] Planned", "Sprint 2+ — pixel obs from fluoroscopy frames", "—"],
        ["Beam hardening (polyenergetic)", "[PLAN] Planned", "Closes realism gap vs DeepDRR", "—"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(5.4),
        headers,
        rows,
        col_widths=[26, 16, 42, 16],
        body_size=11,
        header_size=13,
    )
    return slide


def slide_solvers(prs):
    slide = make_content_slide(prs, "Catheter Physics — Three Solver Backends")
    headers = ["Solver", "Description", "Multi-Env", "Key Features"]
    rows = [
        ["Production XPBD Cosserat rod", "Production-grade solver", "[OK] Yes (batched)", "Mesh BVH collisions, proximal kinematic control, direct solve"],
        ["Self-contained XPBD solver", "Block-Thomas O(n) direct solve", "[OK] Yes (1 thread/rod)", "Batched kernels, floor restitution, GPU root control, CUDA-graph capture, zero Newton dependency"],
        ["Newton XPBD bridge", "External bridge to Newton's XPBD", "[OK] Multi-env (num_envs restriction removed; uses Newton _BatchedRodWorkspace)", "Wraps Newton's block-tridiagonal JMJT solver (PR #1981); no CUDA-graph capture"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.5),
        headers,
        rows,
        col_widths=[24, 26, 18, 32],
        body_size=14,
        header_size=14,
    )
    add_text(
        slide,
        "All three expose 3D centerline positions in metres, scaled ×1000 → mm for the compositing pipeline.",
        Inches(0.4),
        Inches(5.6),
        SLIDE_W - Inches(0.8),
        Inches(0.6),
        size=15,
        color=GREY,
    )
    return slide


def slide_xpbd_multi_env(prs):
    slide = make_content_slide(
        prs,
        "Self-Contained XPBD — Multi-Env Upgrade",
        subtitle="isaaclab_newton/solvers/xpbd_rod_solver.py — feature parity with Newton's xpbd_rod batched path",
    )

    add_text(
        slide,
        "Batched data + kernels",
        Inches(0.5),
        Inches(1.95),
        Inches(6.0),
        Inches(0.4),
        size=16,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_left = [
        "_BatchedWorkspace: concatenated GPU buffers for E rods (positions / orientations / velocities / Jacobians / blocks / λ)",
        "11 batched kernel variants: predict_pos/rot, integrate_pos/rot, prepare_compliance, update_constraints, compute_jacobians, inv_inertia, assemble_jmjt, block_thomas, compute_corrections",
        "Block-Thomas direct solve: 1 GPU thread per rod → cost scales with edges_per_rod, not num_envs × edges",
        "Per-env material parameters: Young's modulus, torsion modulus, gravity, inverse-inertia diag",
        "Index arrays: rod_offsets / edge_offsets / particle_rod_id / edge_rod_id (built once at construction)",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_left],
        Inches(0.5),
        Inches(2.35),
        Inches(6.3),
        Inches(4.7),
        size=12,
        line_spacing=1.3,
    )

    add_text(
        slide,
        "GPU-side control + diagnostics",
        Inches(7.0),
        Inches(1.95),
        Inches(6.0),
        Inches(0.4),
        size=16,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_right = [
        "apply_proximal_control_gpu(push_v[:E], rotate_v[:E], dt) — single Warp launch, no CPU sync",
        "set_root_orientation(env_idx, q) — single-thread launch; safe inside CUDA graph",
        "Floor collision now supports restitution (0 = inelastic, 1 = perfectly elastic)",
        "step() captures the entire substep loop into a CUDA graph the first call; replays via wp.capture_launch (~zero CPU overhead)",
        "max_delta_lambda diagnostic written via wp.atomic_max for convergence monitoring across all envs",
        "Backward-compatible single-env path preserved (~1,300 Hz, 20 segments, A6000)",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_right],
        Inches(7.0),
        Inches(2.35),
        Inches(6.0),
        Inches(4.7),
        size=12,
        line_spacing=1.3,
    )

    add_text(
        slide,
        "Smoke-tested: num_envs=8 step + apply_proximal_control_gpu + set_root_orientation all pass; "
        "positions tensor returns (num_envs, num_points, 3) when batched, (num_points, 3) when single-env.",
        Inches(0.4),
        Inches(7.05),
        SLIDE_W - Inches(0.8),
        Inches(0.5),
        size=11,
        color=DONE_GREEN,
        bold=True,
    )
    return slide


def slide_compositing_paths(prs):
    slide = make_content_slide(prs, "Fluoroscopy Rendering — Three Compositing Paths")
    headers = ["Path", "Compositing", "Performance"]
    rows = [
        ["Slang GPU (unified loop)", "Fused volume + catheter Beer-Lambert in single GPU ray march", "~25 FPS @ 512×512 (A6000)"],
        ["CPU Beer-Lambert", "Per-segment cylinder chord → attenuation map → I = I_DRR × exp(−atten) + scatter + PSF + Poisson", "~2–5 FPS (NumPy)"],
        ["Isaac Lab 3D USD quad", "DRR backdrop + capsule markers in Omniverse viewport", "Real-time (RTX)"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.2),
        headers,
        rows,
        col_widths=[22, 50, 28],
        body_size=14,
        header_size=14,
    )
    add_text(
        slide,
        "Key property: Beer-Lambert is multiplicative — catheter attenuates rather than occludes, with smooth sub-pixel cylinder edges and depth-correct self-crossings.",
        Inches(0.4),
        Inches(5.3),
        SLIDE_W - Inches(0.8),
        Inches(1.0),
        size=15,
        color=NV_GREEN_DARK,
        bold=True,
    )
    return slide


def slide_detector_chain(prs):
    slide = make_content_slide(prs, "Beer-Lambert Detector Physics Chain (CPU path)")
    headers = ["#", "Step", "Purpose"]
    rows = [
        ["1", "Attenuation map: t(d) = 2√(r² − d²) per segment", "Exact cylinder chord thickness"],
        ["2", "I_final = I_DRR × exp(−atten_map)", "Beer-Lambert transmission"],
        ["3", "Veiling glare: σ=18 px Gaussian, 3% re-add", "X-ray scatter in tissue + housing"],
        ["4", "Detector PSF: σ=0.7 px Gaussian", "CsI scintillator finite resolution"],
        ["5", "Poisson noise @ 2000 photons/px", "Low-dose pulsed fluoroscopy statistics"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.6),
        headers,
        rows,
        col_widths=[6, 50, 44],
        body_size=14,
        header_size=14,
    )
    add_text(
        slide,
        "Per-segment 5-zone profile: tungsten markers (μ=3.0) → nitinol braid (μ=0.8) → transition → soft polymer (μ=0.15) → platinum tip (μ=5.0).",
        Inches(0.4),
        Inches(5.7),
        SLIDE_W - Inches(0.8),
        Inches(1.0),
        size=14,
        color=GREY,
    )
    return slide


def slide_xray_perf(prs):
    slide = make_content_slide(prs, "X-Ray Performance Baseline")
    headers = ["Metric", "Target", "Current", "Notes"]
    rows = [
        ["Single-env physics FPS", ">1,000 Hz", "~1,300 Hz (20 seg, A6000)", "[OK] Achieved"],
        ["Slang GPU compositing (fused)", "<5 ms @ 512²", "~40 ms (~25 FPS)", "Single GPU ray march"],
        ["CPU Beer-Lambert compositing", "<2 ms/frame", "~200–500 ms/frame", "NumPy; Warp port → 30+ FPS"],
        ["Multi-env physics (batched XPBD)", ">60 Hz @ 512 envs", "[OK] Available", "Batched block-Thomas + CUDA-graph capture; all 3 backends"],
        ["Multi-env Slang rendering (batched)", ">60 Hz @ 512 envs", "[PARTIAL] ~25 FPS @ N=1; similar at N≤4", "RWStructuredBuffer path live; Texture2DArray upgrade for N>8 = Sprint 2"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.6),
        headers,
        rows,
        col_widths=[28, 18, 28, 26],
        body_size=14,
        header_size=14,
    )
    return slide


def slide_ultrasound(prs):
    slide = make_content_slide(prs, "Ultrasound — What's Implemented")
    bullets = [
        "Slang ultrasound renderer: fixed-step volumetric ray-march + BVH-accelerated triangle intersection",
        "Beer-Lambert attenuation, multi-bounce specular reflection (Phase 2)",
        "Differentiable rendering for 6-DOF probe pose, full PyTorch autograd integration",
        "Legacy OptiX (C++/CUDA) renderer retained as reference (non-differentiable)",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets],
        Inches(0.5),
        Inches(1.7),
        SLIDE_W - Inches(1.0),
        Inches(2.0),
        size=16,
    )

    headers = ["Metric", "Target", "Current"]
    rows = [
        ["Single-env render (Phase 1)", "<5 ms (128 elem, 1024 depth)", "~3–8 ms (est.)"],
        ["Post-processing (PSF + TGC + Hilbert + scan convert)", "<2 ms", "~1–3 ms (GPU PyTorch)"],
        ["End-to-end frame", "<10 ms (100 Hz)", "~8–20 ms (est.)"],
        ["Backward pass (6 DOF)", "<5 ms", "~1–2 ms"],
    ]
    add_table(
        slide,
        Inches(0.5),
        Inches(4.0),
        SLIDE_W - Inches(1.0),
        Inches(2.8),
        headers,
        rows,
        col_widths=[44, 28, 28],
        body_size=13,
        header_size=14,
    )
    return slide


def slide_rl_pipeline(prs):
    slide = make_content_slide(prs, "RL Training Pipeline (State-Based)")
    headers = ["Component", "Description"]
    rows = [
        ["Environment", "Multi-env production solver, proximal push/rotate control, distance-to-target reward"],
        ["RSL-RL wrapper", "Standard VecEnv adapter for RSL-RL PPO"],
        ["PPO config", "Tuned hyperparameters for catheter navigation"],
        ["Training entry", "512 parallel environments, 1500 max iterations"],
        ["Smoke test", "Validates environment without RL dependencies"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.6),
        headers,
        rows,
        col_widths=[26, 74],
        body_size=14,
        header_size=14,
    )
    add_text(
        slide,
        "State-based only today. Pixel/fluoroscopy observations land in Sprint 2+.",
        Inches(0.4),
        Inches(5.7),
        SLIDE_W - Inches(0.8),
        Inches(0.6),
        size=15,
        color=PARTIAL_ORANGE,
        bold=True,
    )
    return slide


def slide_sprint1_completed(prs):
    slide = make_content_slide(
        prs,
        "This Release — Completed Deliverables",
        subtitle="Across IsaacLab catheter solver + i4h-sensor-simulation-internal fluoro-simulator",
    )
    headers = ["Deliverable", "Status / Location"]
    rows = [
        ["Beer-Lambert compositing (physically-correct transmission)", "[OK] CPU + Slang GPU paths"],
        ["Slang GPU fused DRR + catheter single ray march (depth-correct)", "[OK] diffdrr_slang.slang"],
        ["Per-segment attenuation profile (5-zone catheter model)", "[OK] Implemented"],
        ["Cone-beam magnification of projected catheter radius", "[OK] Implemented"],
        ["Detector physics: Poisson + scatter + PSF + gamma", "[OK] realism.py apply_realism()"],
        ["Misregistration jitter (sub-pixel patient motion)", "[OK] realism.apply_misregistration"],
        ["DSA pipeline (4-step: mask + contrast + subtract + post-process)", "[OK] dsa.py DSAPipeline"],
        ["Vessel boost (μ × A=8) on vessel-masked voxels", "[OK] vasculature.apply_vessel_boost"],
        ["VMTK centerline extraction + Dijkstra arrival map", "[OK] vasculature.py"],
        ["Gamma-variate bolus model + per-frame contrast volume", "[OK] vasculature.gamma_variate / build_contrast_volume"],
        ["Per-frame μ update in cine rendering (volume_callback)", "[OK] simulator.render_cine"],
        ["Volumetric instrument injection (max-attenuation, GPU)", "[OK] instrument-injection (Warp atomic_max)"],
        ["9 vendor C-arm presets (GE, Siemens, Philips, Ziehm)", "[OK] CarmGeometry classmethods"],
        ["Differentiable rendering (Slang autodiff, 6-DOF pose gradients)", "[OK] renderDRR_backward kernel"],
        ["Proximal kinematic control API (push / rotate)", "[OK] Implemented"],
        ["Self-contained XPBD: multi-env batched kernels (1 GPU thread / rod)", "[OK] xpbd_rod_solver.py _BatchedWorkspace"],
        ["GPU-side root control: apply_proximal_control_gpu + set_root_orientation", "[OK] CUDA-graph capturable, no CPU sync"],
        ["CUDA-graph capture for the substep loop", "[OK] Auto-captured / replayed in step()"],
        ["Floor collision restitution (0..1)", "[OK] Parity with Newton upstream"],
        ["RL state-based training pipeline (PPO @ 512 envs)", "[OK] Implemented"],
        ["Multi-env Slang fluoroscopy renderer (batched GPU dispatch)", "[OK] renderDRR_forward_batched, StructuredBuffer<Pose/CatheterSegment>"],
        ["Newton XPBD wrapper: multi-env unlocked (num_envs restriction removed)", "[OK] newton_xpbd_rod_wrapper.py — uses Newton _BatchedRodWorkspace"],
        ["E2E workflow validated: catheter motion visible in Beer-Lambert + Slang renders", "[OK] 35 mm traversal, correct Beer-Lambert polarity, NiTi μ_profile"],
        ["XPBDRodSolver: pre/post-constraint hooks (_pre/_post_constraints_hook)", "[OK] Extension points added to _substep + _substep_batched"],
        ["XCathRodSolver: vessel mesh collision (SDF BVH path)", "[OK] xcath_rod_solver.py — _project_vessel_containment_kernel (wp.mesh_query_point_sign_normal)"],
        ["XCathRodSolver: AABB/edge collision path (vertex + segment vs mesh-edge)", "[OK] _project_mesh_vertex/edge_collision_kernel_averaged (wp.mesh_query_aabb)"],
        ["XCathRodSolver: track-guided insertion + configurable collision ordering", "[OK] _track_sliding_kernel; pre/post-constraint stage; Gauss-Seidel iterations"],
        ["E2E vessel collision validation: catheter constrained in synthetic aortic tube", "[OK] docs/e2e_vessel/ — tip follows vessel Z-undulation under gravity + push control"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.05),
        SLIDE_W - Inches(0.8),
        Inches(5.0),
        headers,
        rows,
        col_widths=[60, 40],
        body_size=11,
        header_size=13,
    )
    return slide


def slide_closed_gaps(prs):
    slide = make_content_slide(prs, "This Release — Closed Gaps vs XCATH Requirements")
    headers = ["Capability", "Before", "This Release"]
    rows = [
        ["Catheter Beer-Lambert compositing", "Missing", "[OK] Fused GPU + CPU path"],
        ["Per-segment attenuation profile (5-zone)", "Missing", "[OK] tungsten / nitinol / polymer / platinum"],
        ["Detector physics (Poisson + scatter + PSF)", "Missing", "[OK] realism.py apply_realism()"],
        ["DSA pipeline (4-step: contrast / mask / subtract / post)", "Missing", "[OK] dsa.py DSAPipeline"],
        ["Vessel boost (μ × A=8)", "Missing", "[OK] vasculature.apply_vessel_boost"],
        ["Bolus dynamics — gamma-variate C(t) + per-frame μ", "Missing", "[OK] build_contrast_volume + render_cine"],
        ["VMTK centerline + Dijkstra arrival map", "Missing", "[OK] extract_centerlines + compute_arrival_map"],
        ["Gamma correction (γ=0.8 clinical display TF)", "Missing", "[OK] In RealismConfig + DSASettings"],
        ["Physical scatter convolution", "Partial (2D veiling only)", "[OK] realism.apply_scatter (Gaussian)"],
        ["Misregistration jitter (DSA mask motion)", "Missing", "[OK] realism.apply_misregistration"],
        ["C-arm clinical presets", "Partial (docstrings only)", "[OK] 9 vendor classmethod factories"],
        ["Cone-beam magnification", "Missing", "[OK] Implemented"],
        ["Per-frame μ update in cine rendering", "Missing", "[OK] render_cine(volume_callback=…)"],
        ["Max-attenuation volumetric instrument injection", "Planned", "[OK] Warp atomic_max kernels"],
        ["Fused GPU DRR + catheter (single ray march)", "N/A", "[OK] NEW capability"],
        ["State-based RL pipeline", "Missing", "[OK] PPO @ 512 envs"],
        ["Multi-env XPBD (self-contained solver)", "Single env only", "[OK] Batched block-Thomas, 1 GPU thread/rod"],
        ["GPU-side proximal control (no CPU sync)", "CPU-bridged via PyTorch", "[OK] apply_proximal_control_gpu kernel"],
        ["CUDA-graph capture of substep loop", "Per-step launch overhead", "[OK] Auto-captured in step()"],
        ["Multi-env Slang fluoroscopy renderer", "Single-env only", "[OK] renderDRR_forward_batched (dispatchThreadID.z)"],
        ["Newton XPBD wrapper multi-env", "Single-env restriction", "[OK] num_envs restriction removed; uses _BatchedRodWorkspace"],
        ["E2E catheter motion in renders", "Not validated", "[OK] 35 mm traversal, correct polarity, NiTi μ_profile"],
        ["Vessel mesh collision (catheter-in-vessel containment)", "Not implemented", "[OK] XCathRodSolver: SDF BVH + AABB/edge + track guidance (xcath_rod_solver.py)"],
        ["XPBD pre/post-constraint extension hooks", "Single _project_predicted_positions hook", "[OK] _pre/_post_constraints_hook in _substep + _substep_batched"],
        ["E2E vessel collision validation", "No anatomical constraint", "[OK] Catheter tip follows vessel Z-undulation against gravity (docs/e2e_vessel/)"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(5.4),
        headers,
        rows,
        col_widths=[42, 22, 36],
        body_size=10,
        header_size=13,
    )
    return slide


def slide_carry_overs(prs):
    slide = make_content_slide(
        prs,
        "This Release — Genuinely Remaining Items",
        subtitle="Items still missing after auditing IsaacLab + i4h-sensor-simulation-internal",
    )
    headers = ["Deliverable", "Status", "Notes"]
    rows = [
        ["Image-based RL observations (pixel obs)", "[PLAN] Not wired", "Catheter env is state-only today"],
        ["GPU-side detector physics (scatter / PSF / Poisson on Slang path)", "[PARTIAL] CPU only", "Fine for offline data gen; needed for in-loop GPU training"],
        ["Texture2DArray upgrade for Slang multi-env path", "[PARTIAL] RWStructuredBuffer only", "Needed for cache-efficient rendering at N>8 envs"],
        ["Beam hardening (polyenergetic correction)", "[PLAN] Not implemented", "Monoenergetic only — closes gap vs DeepDRR"],
        ["Selective injection (hemisphere masking)", "[PLAN] Not implemented", "Bilateral CTA → unilateral DSA realism gap"],
        ["Realism evaluation metrics (FID / SSIM / vessel visibility)", "[PLAN] Not implemented", "Needed for agent's iterative refinement loop"],
        ["3D physics-based scatter (vs 2D Gaussian convolution)", "[PARTIAL] 2D model only", "Gaussian scatter approximation in place"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.1),
        SLIDE_W - Inches(0.8),
        Inches(4.6),
        headers,
        rows,
        col_widths=[36, 24, 40],
        body_size=12,
        header_size=14,
    )
    return slide


def slide_next_xray(prs):
    slide = make_content_slide(prs, "Next Release — X-Ray Sprint 2 (Weeks 3–4)", subtitle="Theme: Image observations & GPU detector physics")
    bullets = [
        "[OK] Multi-env Slang renderer COMPLETED — batched renderDRR_forward_batched dispatch is live",
        "[OK] Vessel mesh collision COMPLETED — XCathRodSolver: SDF BVH + AABB/edge path + track-guided insertion",
        "Texture2DArray upgrade for Slang path — improves texture cache hit rate at N>8 envs",
        "Multi-env vessel collision — extend XCathRodSolver to _substep_batched (one mesh per env)",
        "Image-based RL observations — wire multi-env fluoroscopy frames into RL observation dict",
        "GPU-side detector physics — port scatter / PSF / Poisson from NumPy realism.py to Slang/Warp",
        "End-to-end multi-env benchmark: physics + rendering + RL @ 512 envs; target >60 Hz",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets],
        Inches(0.6),
        Inches(2.1),
        SLIDE_W - Inches(1.2),
        Inches(4.2),
        size=18,
        line_spacing=1.35,
    )
    add_text(
        slide,
        "Status today: all three XPBD backends are multi-env. Batched Slang renderer is live. Vessel mesh collision (single-env) "
        "is now implemented via XCathRodSolver. Remaining Sprint 2 items: Texture2DArray, multi-env vessel collision, "
        "image-based RL observations, and GPU-side detector physics.",
        Inches(0.6),
        Inches(6.2),
        SLIDE_W - Inches(1.2),
        Inches(1.0),
        size=13,
        color=GREY,
    )
    return slide


def slide_renderer_multi_env(prs):
    slide = make_content_slide(
        prs,
        "Multi-Env Slang Renderer — COMPLETED",
        subtitle="diffdrr_slang.slang + diffdrr_slang_renderer.py — batched GPU dispatch now live",
    )

    add_text(
        slide,
        "Slang shader changes (diffdrr_slang.slang)",
        Inches(0.5),
        Inches(1.95),
        Inches(6.0),
        Inches(0.4),
        size=15,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_shader = [
        "[OK] renderDRR_forward_batched: [CudaKernel] uint3 dispatchThreadID (z = envIdx)",
        "[OK] StructuredBuffer<Pose> poses — one pose per env; indexed by dispatchThreadID.z",
        "[OK] StructuredBuffer<CatheterSegment> + StructuredBuffer<int> offsets + counts — concatenated batch",
        "[OK] catheterAttenuation_slice / rayMarchWithCatheterSlice helpers for per-env buffer slicing",
        "[OK] Output: flat RWStructuredBuffer<float> at envIdx × H × W + y × W + x",
        "Known limit: RWStructuredBuffer<float> vs RWTexture2D — texture cache miss at large N; Texture2DArray upgrade = Sprint 2",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_shader],
        Inches(0.5),
        Inches(2.35),
        Inches(6.3),
        Inches(4.6),
        size=12,
        line_spacing=1.3,
    )

    add_text(
        slide,
        "Python wrapper (diffdrr_slang_renderer.py)",
        Inches(7.0),
        Inches(1.95),
        Inches(6.0),
        Inches(0.4),
        size=15,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_py = [
        "[OK] SlangDiffDRRRenderer(num_envs=N) — backward-compatible (num_envs=1 unchanged)",
        "[OK] render_batch(poses) + render_batch_with_catheter(poses, catheters) APIs",
        "[OK] Lazy persistent GPU buffers: _ensure_batch_buffers() — no realloc when N is stable",
        "[OK] _upload_poses / _upload_catheters: in-place copy_from_numpy (no new allocation per frame)",
        "[OK] _readback_batch: to_numpy().view(np.float32) — correct typed readback from RWStructuredBuffer",
        "[OK] Smoke-tested: 15/15 tests pass across 1→32 envs, with and without catheter, mixed batches",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_py],
        Inches(7.0),
        Inches(2.35),
        Inches(6.0),
        Inches(4.6),
        size=12,
        line_spacing=1.3,
    )

    add_text(
        slide,
        "[OK] COMPLETED — batched Slang multi-env rendering is live and smoke-tested. "
        "Performance parity with serial path at N≤4; Texture2DArray upgrade improves cache efficiency at N>8 (Sprint 2).",
        Inches(0.4),
        Inches(7.05),
        SLIDE_W - Inches(0.8),
        Inches(0.5),
        size=11,
        color=DONE_GREEN,
        bold=True,
    )
    return slide


def slide_e2e_validation(prs):
    slide = make_content_slide(
        prs,
        "End-to-End Workflow Validation",
        subtitle="Physics + Beer-Lambert compositing + Slang GPU fused render — all stages confirmed",
    )

    add_text(
        slide,
        "Pipeline stages validated",
        Inches(0.5),
        Inches(1.95),
        Inches(6.3),
        Inches(0.4),
        size=15,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_stages = [
        "[OK] Stage 2: Production RodSolver + CPU Beer-Lambert AP view — catheter advances 35 mm across cranial anatomy",
        "[OK] Stage 3: XPBDRodSolver (self-contained, CUDA graphs) + Beer-Lambert — tip traverses 35 mm in 1 s simulation",
        "[OK] Stage 4: Production RodSolver dual AP + Lateral views — same catheter trajectory visible in both projections",
        "[OK] Stage 5: Multi-env XPBD (4 envs, separate solvers) — tip position traces diverge as push velocity differs",
        "[OK] Stage 6: Pure Slang DRR at 4 C-arm angles (AP / LAO 30° / Lateral / RAO 30°) with real cranial CT",
        "[OK] Stage 7: Slang GPU fused volume + XPBD catheter — catheter enters right FOV edge, advances 34 mm leftward",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_stages],
        Inches(0.5),
        Inches(2.35),
        Inches(6.3),
        Inches(4.2),
        size=12,
        line_spacing=1.35,
    )

    add_text(
        slide,
        "Key fixes applied during validation",
        Inches(7.0),
        Inches(1.95),
        Inches(6.0),
        Inches(0.4),
        size=15,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_fixes = [
        "[OK] Beer-Lambert polarity: I_final = 1 − (1−I_bg) × t_catheter — catheter now bright white (radio-opaque)",
        "[OK] μ_profile calibrated: NiTi shaft 0.40 mm⁻¹ tapering to 0.03 mm⁻¹ at soft polymer tip",
        "[OK] Stage 7 centering: offset computed once at frame 0 — catheter moves through volume, not snaps back",
        "[OK] Duplicate PUSH_VEL constant found and removed — proximal control was silently overridden to 8 mm/s",
        "[OK] Catheter enters from right FOV edge at t=0; traverses anatomy leftward at 40 mm/s proximal push",
        "[OK] 30° Y-rotation applied consistently across Beer-Lambert + Slang stages for matching projection geometry",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_fixes],
        Inches(7.0),
        Inches(2.35),
        Inches(6.0),
        Inches(4.2),
        size=12,
        line_spacing=1.35,
    )

    add_text(
        slide,
        "Output images: source/isaaclab_newton/docs/e2e_output_v2/  (s2–s7 stages, 22 frames total)",
        Inches(0.4),
        Inches(7.05),
        SLIDE_W - Inches(0.8),
        Inches(0.45),
        size=11,
        color=GREY,
        bold=False,
    )
    return slide


def slide_next_us(prs):
    slide = make_content_slide(prs, "Next Release — Ultrasound Sprint 1 (Isaac Lab Integration)")
    headers = ["Deliverable", "Description"]
    rows = [
        ["Ultrasound sensor wrapper", "First-class Isaac Lab sensor module"],
        ["CT-to-volume pipeline", "HU → acoustic impedance / scattering coefficient mapping"],
        ["Observation dict integration", "RL-ready B-mode frames as observations"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.0),
        headers,
        rows,
        col_widths=[32, 68],
        body_size=15,
        header_size=14,
    )
    return slide


def slide_next_phase1(prs):
    slide = make_content_slide(
        prs,
        "Next Release — Phase 1 Fidelity Items (X-Ray)",
        subtitle="Remaining simulation-fidelity work after current implementation audit",
    )
    headers = ["Effort", "Deliverable"]
    rows = [
        ["[OK] DONE", "Multi-env / batched Slang rendering — renderDRR_forward_batched + StructuredBuffer batch path"],
        ["[OK] DONE", "Vessel mesh collision — XCathRodSolver (SDF BVH + AABB/edge + track-guided insertion, single-env)"],
        ["~2 days", "Texture2DArray upgrade — replace RWStructuredBuffer output with hardware-cached array for N>8 envs"],
        ["~2 days", "Multi-env vessel collision — extend XCathRodSolver._pre/_post hooks to _substep_batched (one mesh per env)"],
        ["~3 days", "Image-based RL observations — pixel obs from fluoroscopy frames into obs dict"],
        ["~3 days", "GPU-side detector physics on Slang path (scatter / PSF / Poisson kernels)"],
        ["~2 days", "Beam hardening (polyenergetic correction)"],
        ["~1 day", "Selective injection (hemisphere masking) on top of bolus pipeline"],
        ["~2 days", "Realism evaluation module (FID, SSIM, vessel visibility metrics)"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.1),
        SLIDE_W - Inches(0.8),
        Inches(3.8),
        headers,
        rows,
        col_widths=[14, 86],
        body_size=14,
        header_size=14,
    )
    add_text(
        slide,
        "DSA, vessel boost, bolus dynamics, gamma, scatter, jitter, C-arm presets, and per-frame μ updates are already implemented "
        "in i4h-sensor-simulation-internal — the remaining work focuses on multi-env scaling and image-based RL.",
        Inches(0.4),
        Inches(6.1),
        SLIDE_W - Inches(0.8),
        Inches(1.0),
        size=12,
        color=NV_GREEN_DARK,
        bold=True,
    )
    return slide


def slide_sprint3(prs):
    slide = make_content_slide(prs, "Sprint 3 — Training Readiness (Weeks 5–6)", subtitle="X-Ray catheter navigation policy training")
    bullets = [
        "Domain randomization — C-arm angles, attenuation, photon count, scatter parameters",
        "Gymnasium wrapper — standard RL ecosystem compatibility",
        "Per-task CUDA-graph variants — separate graphs for sub-batches with mixed materials",
        "Image-based observations — fluoroscopy frames as RL inputs",
        "Automated pytest suite — regression coverage for solver + compositing + renderer",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets],
        Inches(0.6),
        Inches(2.1),
        SLIDE_W - Inches(1.2),
        Inches(4.5),
        size=20,
        line_spacing=1.5,
    )
    return slide


def slide_skill_packaging(prs):
    slide = make_content_slide(prs, "Agentic Workflow — Phase 2 (Skill Packaging)", subtitle="Each pipeline stage wrapped as a portable OpenClaw / NemoClaw Skill")
    headers = ["Skill", "Purpose"]
    rows = [
        ["patient-digital-twin", "CTA → μ volume + vessel mask + centerline + arrival map"],
        ["catheter-physics-sim", "Newton rod solver + compositing config"],
        ["sensor-sim-xray", "DRR / DSA / vessel-boost rendering modes"],
        ["dataset-creation", "Paired HDF5/WebDataset (frames + pose + GT)"],
        ["reward-function", "RL reward configuration (target, contact, dose, progress, success)"],
        ["policy-training", "IL (GR00T-H) → RL (PPO/SAC) → SIL evaluation"],
        ["evaluation", "Success rate, navigation time, contact force, FID, registration accuracy"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.1),
        SLIDE_W - Inches(0.8),
        Inches(4.8),
        headers,
        rows,
        col_widths=[26, 74],
        body_size=14,
        header_size=14,
    )
    return slide


def slide_agent_integration(prs):
    slide = make_content_slide(prs, "Agentic Workflow — Phase 3 (Agent Integration)")
    headers = ["Week", "Deliverable"]
    rows = [
        ["5–6", "Skill discovery & chaining — parse skill definitions, resolve I/O dependencies"],
        ["6", 'Natural language → config mapping — "use Philips Azurion, DSA mode" → YAML overrides'],
        ["6–7", "Iterative refinement loop — agent runs evaluation, analyzes metrics, proposes changes, re-runs"],
        ["7", "Slack / IDE integration — agent posts progress, visualizations, and final reports"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(1.7),
        SLIDE_W - Inches(0.8),
        Inches(3.4),
        headers,
        rows,
        col_widths=[10, 90],
        body_size=14,
        header_size=14,
    )
    add_text(
        slide,
        "Outcome: experiment cycle compresses from weeks → hours, with the agent running continuously.",
        Inches(0.4),
        Inches(5.6),
        SLIDE_W - Inches(0.8),
        Inches(0.7),
        size=18,
        color=NV_GREEN_DARK,
        bold=True,
    )
    return slide


def slide_remaining_gaps(prs):
    slide = make_content_slide(
        prs,
        "Following Releases — Remaining Capability Gaps",
        subtitle="Genuinely-missing items after auditing IsaacLab + i4h-sensor-simulation-internal",
    )
    headers = ["Missing Feature", "Owning Skill", "Why It Matters"]
    rows = [
        ["Texture2DArray for Slang multi-env path", "sensor-sim-xray", "Cache-efficient rendering at N>8 envs (RWStructuredBuffer path live)"],
        ["Image-based RL observations", "reward-function, policy-training", "Image-guided policies (env is state-only today)"],
        ["GPU-side detector physics on Slang path", "sensor-sim-xray", "In-loop GPU training without CPU readback"],
        ["Beam hardening (polyenergetic)", "sensor-sim-xray", "Closes realism gap vs DeepDRR"],
        ["Selective injection (hemisphere masking)", "patient-digital-twin", "Bilateral CTA → unilateral DSA realism"],
        ["Realism metrics (FID, SSIM, vessel visibility)", "evaluation", "Quantitative loop for iterative refinement"],
        ["3D physics-based scatter (vs 2D Gaussian)", "sensor-sim-xray", "Higher-fidelity scatter halo"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.1),
        SLIDE_W - Inches(0.8),
        Inches(4.8),
        headers,
        rows,
        col_widths=[36, 28, 36],
        body_size=12,
        header_size=14,
    )
    add_text(
        slide,
        "DSA, vessel boost, bolus tracking, per-frame μ update, max-attenuation volume compositing, "
        "C-arm presets, gamma, scatter convolution, misregistration jitter, multi-env Slang renderer, "
        "Newton wrapper multi-env, E2E workflow validation, and vessel mesh collision (XCathRodSolver) — all closed in this release.",
        Inches(0.4),
        Inches(6.05),
        SLIDE_W - Inches(0.8),
        Inches(0.9),
        size=12,
        color=DONE_GREEN,
        bold=True,
    )
    return slide


def slide_future_sensors(prs):
    slide = make_content_slide(
        prs,
        "Workflow Enhancements",
        subtitle="Capabilities that further augment the X-Ray catheter intervention workflow",
    )
    headers = ["Capability", "Priority", "Integration Effort", "Why It Matters for X-Ray Catheter Workflow"]
    rows = [
        ["Texture2DArray upgrade (Slang multi-env)", "High", "~2 days — replace RWStructuredBuffer output", "Improves texture cache hit rate at N>8; batched path is live at N≤4"],
        ["Image-based RL observations", "High", "~3 days — wire pixels → obs dict", "Enables image-guided policies; current env is state-only"],
        ["Force / Torque sensing", "High", "Low — data already in collision solver", "Safety reward signal; tip-force penalty for vessel perforation"],
        ["Beam hardening (polyenergetic)", "Medium", "Moderate", "Closes realism gap vs DeepDRR; improves sim-to-real transfer"],
        ["GPU-side detector physics on Slang path", "Medium", "~3 days — Warp/Slang kernels", "Avoids CPU readback in the in-loop training pipeline"],
        ["FID / SSIM / vessel-visibility realism metrics", "Medium", "~2 days", "Closes the agent's iterative refinement loop"],
        ["CBCT reconstruction", "Low–Medium", "Moderate — batched DRR + FDK", "Intra-procedural 3D imaging; reuses existing GPU ray-caster"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.1),
        SLIDE_W - Inches(0.8),
        Inches(4.6),
        headers,
        rows,
        col_widths=[26, 12, 26, 36],
        body_size=11,
        header_size=13,
    )
    return slide


def slide_dsa_bolus(prs):
    slide = make_content_slide(
        prs,
        "DSA Pipeline & Bolus Dynamics — Implemented",
        subtitle="fluorosim/dsa.py + fluorosim/vasculature.py",
    )

    add_text(
        slide,
        "DSA Pipeline (DSAPipeline.render_dsa_frame)",
        Inches(0.5),
        Inches(2.0),
        Inches(6.0),
        Inches(0.4),
        size=16,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_dsa = [
        "Step 1–2: render mask DRR + contrast DRR at same pose",
        "Step 3: scatter convolution applied to both (Gaussian σ_s)",
        "Step 4: independent Poisson + Gaussian noise per DRR",
        "Step 5: misregistration jitter (sub-pixel shift) on mask",
        "Step 6: subtract  diff = contrast − mask",
        "Step 7: contrast boost (k=20) → gamma (γ=0.8) → normalize",
        "Plus render_dsa_sequence() for temporal cine output",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_dsa],
        Inches(0.5),
        Inches(2.4),
        Inches(6.2),
        Inches(4.5),
        size=13,
        line_spacing=1.25,
    )

    add_text(
        slide,
        "Bolus Dynamics (vasculature.py)",
        Inches(7.0),
        Inches(2.0),
        Inches(6.0),
        Inches(0.4),
        size=16,
        bold=True,
        color=NV_GREEN_DARK,
    )
    bullets_bolus = [
        "extract_centerlines: VMTK marching-cubes + network extraction",
        "compute_arrival_map: Dijkstra travel-time over centerline graph",
        "gamma_variate(t, α, β): C(t) = (t/t_peak)^α · exp(α(1−t/t_peak))",
        "build_contrast_volume: μ(v,t) = μ_tissue + Δμ · C(t − T(v))",
        "apply_vessel_boost: μ × A=8 on vessel-masked voxels",
        "Wired into FluoroSimulator.render_cine via volume_callback",
        "Enables 150-frame @ 5 FPS = 30 s temporal DSA simulation",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_bolus],
        Inches(7.0),
        Inches(2.4),
        Inches(6.0),
        Inches(4.5),
        size=13,
        line_spacing=1.25,
    )
    return slide


def slide_instrument_injection(prs):
    slide = make_content_slide(
        prs,
        "Volumetric Instrument Injection — Implemented",
        subtitle="instrument-injection module — max-attenuation compositing on GPU",
    )

    add_text(
        slide,
        "Compositing Rule",
        Inches(0.5),
        Inches(2.0),
        Inches(6.0),
        Inches(0.4),
        size=16,
        bold=True,
        color=NV_GREEN_DARK,
    )
    add_text(
        slide,
        "μ_composited(v) = max(μ_anatomy(v), μ_instrument)",
        Inches(0.5),
        Inches(2.4),
        Inches(6.5),
        Inches(0.5),
        size=15,
        bold=True,
        color=BLACK,
        font="Consolas",
    )
    bullets_a = [
        "Bone preserved: μ_bone > μ_polymer → catheter behind bone",
        "Soft tissue replaced: μ_catheter > μ_tissue → catheter visible",
        "Multi-instrument: densest material wins per voxel",
        "Implemented as wp.atomic_max — thread-safe parallel compositing",
    ]
    add_bullets(
        slide,
        [(b, False) for b in bullets_a],
        Inches(0.5),
        Inches(3.0),
        Inches(6.2),
        Inches(2.5),
        size=13,
        line_spacing=1.3,
    )

    add_text(
        slide,
        "Instruments + Performance",
        Inches(7.0),
        Inches(2.0),
        Inches(6.0),
        Inches(0.4),
        size=16,
        bold=True,
        color=NV_GREEN_DARK,
    )
    headers = ["Material", "μ (mm⁻¹)", "Use"]
    rows = [
        ["Polymer", "0.03", "Catheter shaft"],
        ["Nitinol", "0.08", "Guidewire / shaft"],
        ["Stainless steel", "0.10", "Needle"],
        ["Platinum", "0.20", "Marker bands"],
        ["Tungsten", "0.35", "Radiopaque markers"],
        ["Iodine contrast", "0.05", "DSA contrast agent"],
    ]
    add_table(
        slide,
        Inches(7.0),
        Inches(2.45),
        Inches(5.9),
        Inches(2.4),
        headers,
        rows,
        col_widths=[40, 22, 38],
        body_size=11,
        header_size=12,
    )
    add_text(
        slide,
        "Kernels: paint_spheres_kernel + paint_cylinders_kernel (Warp)\n"
        "Performance: ~2 ms / 64 catheter nodes / 512³ volume\n"
        "(vs ~500 ms for the equivalent CPU painting loop)",
        Inches(7.0),
        Inches(5.2),
        Inches(5.9),
        Inches(1.5),
        size=12,
        color=GREY,
    )
    return slide


def slide_carm_presets(prs):
    slide = make_content_slide(
        prs,
        "C-arm Vendor Presets — Implemented",
        subtitle="CarmGeometry classmethod factories — fluorosim/config.py",
    )
    headers = ["Vendor / Model", "SDD (mm)", "SID (mm)", "Detector", "Pixel (mm)"]
    rows = [
        ["GE OEC 9900", "1020", "510", "1024 × 1024", "0.194"],
        ["GE OEC Elite CFD", "1150", "575", "1920 × 1920", "0.154"],
        ["GE Innova IGS 540", "1200", "750", "2048 × 2048", "0.200"],
        ["Siemens Arcadis Avantic", "1000", "500", "1024 × 1024", "0.195"],
        ["Siemens Cios Alpha", "1100", "550", "1536 × 1536", "0.178"],
        ["Siemens Artis zee", "1250", "780", "2480 × 1920", "0.154"],
        ["Philips BV Pulsera", "990", "495", "1024 × 1024", "0.200"],
        ["Philips Azurion 7", "1240", "780", "2480 × 1920", "0.154"],
        ["Ziehm Vision RFD 3D", "1000", "500", "1024 × 1024", "0.194"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(2.1),
        SLIDE_W - Inches(0.8),
        Inches(4.5),
        headers,
        rows,
        col_widths=[34, 14, 14, 24, 14],
        body_size=12,
        header_size=13,
    )
    add_text(
        slide,
        "Usage:  geometry = CarmGeometry.philips_azurion_7()   →   one-line agent-selectable preset",
        Inches(0.4),
        Inches(6.7),
        SLIDE_W - Inches(0.8),
        Inches(0.5),
        size=13,
        color=NV_GREEN_DARK,
        bold=True,
        font="Consolas",
    )
    return slide


def slide_summary(prs):
    slide = make_content_slide(prs, "Summary")

    col_w = (SLIDE_W - Inches(1.6)) / 4
    col_y = Inches(1.7)
    col_h = Inches(5.4)

    def column(x, label, accent, body):
        add_rect(slide, x, col_y, col_w - Inches(0.15), Inches(0.5), accent)
        add_text(
            slide,
            label,
            x,
            col_y + Inches(0.05),
            col_w - Inches(0.15),
            Inches(0.4),
            size=15,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        tb = slide.shapes.add_textbox(x, col_y + Inches(0.6), col_w - Inches(0.15), col_h - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.3
            p.space_after = Pt(4)
            r = p.add_run()
            r.text = f"• {line}"
            r.font.name = "Calibri"
            r.font.size = Pt(13)
            r.font.color.rgb = BLACK

    x0 = Inches(0.4)
    column(
        x0,
        "TODAY",
        GREY,
        [
            "X-Ray fluoroscopy pipeline implemented end-to-end",
            "Fused GPU Beer-Lambert at ~25 FPS, physics ~1,300 Hz",
            "DSA + bolus + vessel boost + 9 vendor C-arm presets",
            "Volumetric instrument injection (atomic-max, ~2 ms)",
            "State-based PPO @ 512 envs operational",
        ],
    )
    column(
        x0 + col_w,
        "THIS RELEASE",
        DONE_GREEN,
        [
            "Full DSA pipeline (4-step) + temporal bolus dynamics",
            "Detector physics: Poisson + scatter + PSF + gamma + jitter",
            "Self-contained XPBD: batched + GPU control + CUDA graphs",
            "Multi-env Slang renderer: renderDRR_forward_batched live",
            "Newton XPBD wrapper: multi-env unlocked",
            "E2E validated: 35 mm catheter traversal in Slang renders",
            "Vessel mesh collision: XCathRodSolver (SDF BVH + AABB/edge + track)",
            "21+ XCATH-required capabilities CLOSED",
        ],
    )
    column(
        x0 + col_w * 2,
        "NEXT RELEASE",
        PARTIAL_ORANGE,
        [
            "Texture2DArray upgrade (N>8 env cache efficiency)",
            "Multi-env vessel collision (one mesh per env)",
            "Image-based RL observations (pixel obs)",
            "GPU-side detector physics on Slang path",
            "Beam hardening (polyenergetic correction)",
        ],
    )
    column(
        x0 + col_w * 3,
        "FOLLOWING",
        NV_GREEN_DARK,
        [
            "Sprint 3 training readiness (DR, Gymnasium, per-task graphs)",
            "Phase 2 skill packaging (7 OpenClaw skills)",
            "Phase 3 agent integration (NL → config → eval loop)",
            "Realism metrics (FID / SSIM / vessel visibility)",
            "Workflow extensions: F/T, CBCT",
        ],
    )
    return slide


def slide_xcath_rod_solver(prs):
    slide = make_content_slide(
        prs,
        "XCathRodSolver — Vessel Mesh Collision COMPLETED",
        subtitle="xcath_rod_solver.py — ported from Newton xpbd_rods_solver_integr branch (Przemek Korzeniowski)",
    )

    # Left column: architecture
    add_text(
        slide,
        "Architecture",
        Inches(0.5),
        Inches(1.95),
        Inches(6.0),
        Inches(0.4),
        size=14,
        bold=True,
    )
    arch_rows = [
        ("XCathRodSolver(XPBDRodSolver)", "Subclass — no changes to base solver interface"),
        ("_pre_constraints_hook(ws, dt, dev)", "Vessel containment before XPBD solve (optional)"),
        ("_post_constraints_hook(ws, dt, dev)", "Track guidance + vessel containment after XPBD solve"),
        ("collision_mesh: wp.Mesh", "Arbitrary closed triangle mesh (BVH auto-built by Warp)"),
        ("track_start / track_dir / track_length", "Insertion axis — non-tip particles snap to axis"),
        ("tip_num_edges", "Distal free-tip segment count (steerable, unconstrained)"),
        ("collision_iterations", "Gauss-Seidel passes per substep (default 2)"),
        ("sign_scale / target_phi / max_dist", "SDF convention, clearance offset, BVH search radius"),
    ]

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.45), Inches(6.3), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (name, desc) in enumerate(arch_rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = f"{name}  "
        r1.font.name = "Courier New"
        r1.font.size = Pt(10)
        r1.font.bold = True
        r1.font.color.rgb = NV_GREEN_DARK
        r2 = p.add_run()
        r2.text = f"— {desc}"
        r2.font.name = "Calibri"
        r2.font.size = Pt(10)
        r2.font.color.rgb = BLACK

    # Right column: collision paths
    add_text(
        slide,
        "Two Collision Paths",
        Inches(7.0),
        Inches(1.95),
        Inches(5.9),
        Inches(0.4),
        size=14,
        bold=True,
    )
    right_bullets = [
        "SDF path (default): wp.mesh_query_point_sign_normal",
        "  → BVH signed-distance, ~O(log N) per particle",
        "  → Push particle to target_phi clearance from wall",
        "  → Supports smooth vertex normals (area-weighted)",
        "",
        "AABB / edge path (use for tight curvature):",
        "  → _project_mesh_vertex_collision_kernel_averaged",
        "     AABB broadphase, vertex-vs-triangle contacts",
        "  → _project_mesh_edge_collision_kernel_averaged",
        "     Rod-segment vs mesh-edge closest-point contacts",
        "  → Corrections flushed atomically per Gauss-Seidel iter",
        "",
        "E2E validation: synthetic aortic tube (13 mm radius,",
        "  Z-undulating path). Tip tracks vessel Z-curve",
        "  against gravity (Δz = +2.7 mm over 40 mm push). ✓",
    ]
    add_bullets(
        slide,
        [(b, False) for b in right_bullets],
        Inches(7.0),
        Inches(2.45),
        Inches(5.9),
        Inches(4.1),
        size=11,
        line_spacing=1.2,
    )

    add_text(
        slide,
        "Base class XPBDRodSolver now exposes _pre/_post_constraints_hook in both _substep (single-env) "
        "and _substep_batched (multi-env) for future per-env mesh collision at 512+ envs.",
        Inches(0.5),
        Inches(6.55),
        SLIDE_W - Inches(1.0),
        Inches(0.7),
        size=12,
        color=GREY,
    )
    return slide


def slide_questions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.06), NV_GREEN)
    add_text(
        slide,
        "Questions?",
        Inches(0.6),
        Inches(2.5),
        SLIDE_W - Inches(1.2),
        Inches(1.2),
        size=72,
        bold=True,
        color=NV_GREEN,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "X-Ray–Guided Robotic Catheter Intervention — Isaac for Healthcare",
        Inches(0.6),
        Inches(3.6),
        SLIDE_W - Inches(1.2),
        Inches(0.5),
        size=22,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "NVIDIA Healthcare — Holoscan Team",
        Inches(0.6),
        Inches(4.1),
        SLIDE_W - Inches(1.2),
        Inches(0.5),
        size=16,
        color=LIGHT_GREY,
        align=PP_ALIGN.CENTER,
    )
    return slide


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("title", make_title_slide),
        ("agenda", slide_agenda),
        ("xcath_req", slide_xcath_requirements),
        ("xcath_prog", slide_xcath_progress),
        ("exec", slide_executive_snapshot),
        ("modality", slide_modality_status),
        ("part1", lambda p: make_section_divider(p, "PART 1", "Current Status", "What is built, integrated, and measured today.")),
        ("solvers", slide_solvers),
        ("xpbd_multi_env", slide_xpbd_multi_env),
        ("compositing", slide_compositing_paths),
        ("detector", slide_detector_chain),
        ("dsa_bolus", slide_dsa_bolus),
        ("instr_inject", slide_instrument_injection),
        ("carm_presets", slide_carm_presets),
        ("xray_perf", slide_xray_perf),
        ("rl_pipeline", slide_rl_pipeline),
        ("part2", lambda p: make_section_divider(p, "PART 2", "This Release", "Sprint 1 — features shipping now.")),
        ("sprint1_done", slide_sprint1_completed),
        ("closed_gaps", slide_closed_gaps),
        ("carry_overs", slide_carry_overs),
        ("renderer_multi_env", slide_renderer_multi_env),
        ("e2e_validation", slide_e2e_validation),
        ("xcath_solver", slide_xcath_rod_solver),
        ("part3", lambda p: make_section_divider(p, "PART 3", "Next Release", "Sprint 2 — Texture2DArray, multi-env vessel collision, image observations.")),
        ("next_xray", slide_next_xray),
        ("next_phase1", slide_next_phase1),
        ("part4", lambda p: make_section_divider(p, "PART 4", "Following Releases", "Sprint 3 + agentic workflow integration.")),
        ("sprint3", slide_sprint3),
        ("skills", slide_skill_packaging),
        ("agent", slide_agent_integration),
        ("gaps", slide_remaining_gaps),
        ("part5", lambda p: make_section_divider(p, "PART 5", "Workflow Enhancements", "Adjacent capabilities augmenting the X-Ray catheter workflow.")),
        ("future", slide_future_sensors),
        ("summary", slide_summary),
        ("qna", slide_questions),
    ]

    for _, fn in builders:
        fn(prs)

    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1 or idx == total:
            continue
        add_footer(slide, idx, total)

    out = Path(__file__).resolve().parent / "sensor_simulation_release_deck.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    build()
