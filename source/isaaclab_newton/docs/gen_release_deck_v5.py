"""Generate release_deck_v5.pptx — X-Ray–Guided Robotic Catheter Intervention.

New slides vs v4:
  - "Unified Sim Loop — How It Works" (detailed algorithm walkthrough)
  - "Newton Upstream Algorithms" (what was ported from the Newton research branch)
  - Corrected "Fluoroscopy Rendering" slide (accurate FPS, CPU vs GPU paths)
  - Updated performance table (256×256 interactive + 512×512 benchmark)

Run:
    python3 gen_release_deck_v5.py
Output:
    release_deck_v5.pptx  (same directory)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = Path(__file__).parent / "release_deck_v5.pptx"

# ── brand colours ──────────────────────────────────────────────────────────────
NV_GREEN   = RGBColor(0x76, 0xB9, 0x00)   # NVIDIA green
NV_DARK    = RGBColor(0x1A, 0x1A, 0x1A)   # near-black
NV_GREY    = RGBColor(0xF2, 0xF2, 0xF2)   # light grey
NV_MID     = RGBColor(0x55, 0x55, 0x55)   # mid grey
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x00, 0xA0, 0xD0)   # blue accent
VALID_COL  = RGBColor(0xAA, 0x66, 0x00)   # amber — validated, not yet merged
IMPL_COL   = RGBColor(0x16, 0xA0, 0x85)   # teal  — implemented in codebase
ANALY_COL  = RGBColor(0x77, 0x77, 0x77)   # grey  — analysis result only

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=NV_DARK, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle=None):
    """Green top bar with title."""
    add_rect(slide, 0, 0, 13.33, 0.75, fill=NV_GREEN)
    add_text(slide, "X-Ray–Guided Robotic Catheter Intervention — Isaac for Healthcare",
             0.15, 0.03, 13.0, 0.4, size=11, bold=False, color=NV_DARK)
    add_rect(slide, 0, 0.75, 13.33, 0.05, fill=NV_DARK)
    add_text(slide, title, 0.3, 0.85, 12.5, 0.8, size=28, bold=True, color=NV_DARK)
    if subtitle:
        add_text(slide, subtitle, 0.3, 1.55, 12.5, 0.4, size=13, color=NV_MID, italic=True)

def bullet_block(slide, items, l, t, w, h, title=None, title_color=NV_GREEN,
                 item_size=13, item_color=NV_DARK, title_size=15):
    """Draw a titled bullet block."""
    y = t
    if title:
        add_text(slide, title, l, y, w, 0.35, size=title_size, bold=True,
                 color=title_color)
        y += 0.32
    for item in items:
        add_text(slide, f"• {item}", l+0.1, y, w-0.1, 0.28,
                 size=item_size, color=item_color)
        y += 0.27

def table_slide(slide, headers, rows, l, t, w, col_widths=None):
    """Draw a simple table with green header row."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols
    row_h = 0.32
    # Header row
    x = l
    for i, hdr in enumerate(headers):
        add_rect(slide, x, t, col_widths[i], row_h, fill=NV_GREEN)
        add_text(slide, hdr, x+0.05, t+0.02, col_widths[i]-0.1, row_h-0.05,
                 size=12, bold=True, color=NV_DARK)
        x += col_widths[i]
    # Data rows
    for ri, row in enumerate(rows):
        x = l
        fill = NV_GREY if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            add_rect(slide, x, t+(ri+1)*row_h, col_widths[ci], row_h,
                     fill=fill, line=RGBColor(0xCC,0xCC,0xCC))
            add_text(slide, cell, x+0.05, t+(ri+1)*row_h+0.02,
                     col_widths[ci]-0.1, row_h-0.05, size=11, color=NV_DARK)
            x += col_widths[ci]


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NV_DARK)
add_rect(sl, 0, 0, 13.33, 0.12, fill=NV_GREEN)
add_rect(sl, 0, 7.38, 13.33, 0.12, fill=NV_GREEN)
add_text(sl, "X-Ray–Guided Robotic Catheter Intervention",
         0.6, 1.8, 12.0, 1.2, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Simulation Workflow — Release Status & Roadmap",
         0.6, 2.9, 12.0, 0.7, size=24, bold=False, color=NV_GREEN, align=PP_ALIGN.LEFT)
add_rect(sl, 0.6, 3.65, 4.0, 0.04, fill=NV_GREEN)
bullets = [
    "Cosserat rod XPBD physics  (Newton upstream + vessel-mesh collision)",
    "Slang GPU Beer-Lambert DRR renderer  (inline catheter compositing)",
    "Unified sim loop  (physics → render → RL obs, zero CPU sync on hot path)",
    "Interactive UI demo  +  RL training pipeline  (RSL-RL PPO, 512 envs)",
]
y = 3.85
for b in bullets:
    add_text(sl, f"▸  {b}", 0.6, y, 12.0, 0.35, size=15, color=RGBColor(0xCC,0xCC,0xCC))
    y += 0.38
add_text(sl, "Isaac for Healthcare  |  May 2026",
         0.6, 6.9, 12.0, 0.4, size=12, color=NV_MID, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ARCHITECTURE OVERVIEW (4-stage pipeline)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Architecture Overview — Four-Stage Pipeline")

stages = [
    ("Stage 1\nCT Ingestion", "CT/CTA  →  HU→μ conversion  →  3D μ-volume\n"
      "Vessel segmentation  →  wp.Mesh (collision)\n"
      "Centerline (VMTK)  →  Dijkstra arrival map",
      0.25),
    ("Stage 2\nSim Environment", "XCathRodSolver  (XPBD + vessel-mesh SDF)\n"
      "Slang DRR renderer  (Beer-Lambert, inline catheter)\n"
      "Unified loop: physics → render → RL obs",
      3.55),
    ("Stage 3\nRL Training", "RSL-RL PPO  |  512 parallel envs\n"
      "State obs today  →  Pixel obs Sprint 2\n"
      "GR00T-N IL + RL fine-tune",
      6.85),
    ("Stage 4\nDeployment", "Holoscan IGX  |  <10 ms inference\n"
      "XCath robotic catheter  |  Safety layer\n"
      "Sim-to-real via domain randomization",
      10.15),
]
for title, body, x in stages:
    add_rect(sl, x, 2.1, 2.9, 4.6, fill=NV_GREY, line=NV_GREEN)
    add_rect(sl, x, 2.1, 2.9, 0.55, fill=NV_GREEN)
    add_text(sl, title, x+0.1, 2.12, 2.7, 0.5, size=13, bold=True, color=NV_DARK)
    add_text(sl, body, x+0.12, 2.72, 2.68, 3.8, size=11, color=NV_DARK)

# Arrows between stages
for ax in [3.2, 6.5, 9.8]:
    add_text(sl, "▶", ax, 4.1, 0.3, 0.4, size=18, bold=True, color=NV_GREEN)

add_text(sl,
    "All four stages share one coordinate frame (CT-volume mm).  "
    "The same μ-volume fed to the renderer is derived from the same CT as the collision mesh — "
    "physics and imaging are always registered to the same patient anatomy.",
    0.3, 6.85, 12.7, 0.5, size=11, color=NV_MID, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — UNIFIED SIM LOOP — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Unified Sim Loop — Per-Timestep Execution")
add_text(sl,
    "Five GPU-resident steps per frame.  No CPU↔GPU synchronization on the hot path.",
    0.3, 1.6, 12.7, 0.35, size=13, color=NV_MID, italic=True)

steps = [
    ("Step 1", "Root Control",
     "Proximal boundary condition updated via GPU kernel.\n"
     "Insertion depth δ and rotation angle θ written directly to\n"
     "the proximal particle position + orientation buffers.\n"
     "Warp kernels: apply_proximal_control_gpu  /  set_root_orientation\n"
     "Both are CUDA-graph capturable — zero CPU sync."),
    ("Step 2", "Physics Step  (XPBD CUDA graph)",
     "CUDA graph replayed — entire N-substep loop fires as one graph node.\n"
     "Each substep: predict → assemble JMJT → block-Thomas direct solve → apply Δx → integrate.\n"
     "Vessel-mesh SDF + AABB collision constraints applied before/after solve.\n"
     "1,300 Hz throughput per environment.  Scales to 512 envs via flat buffer offsets."),
    ("Step 3", "Segment Buffer Update",
     "Particle positions (N×3 float32, GPU-resident) written to\n"
     "the renderer's CatheterSegmentData buffer (N×8 floats: xyz0, xyz1, radius, μ).\n"
     "No host transfer.  No CT volume modification."),
    ("Step 4", "Render Call  (Slang DRR)",
     "Ray-march kernel dispatched over H×W detector grid.\n"
     "CT μ-volume is a static GPU texture — never reloaded per frame.\n"
     "Catheter composited inline: μ_total(s) = μ_CT(s) + Σᵢ μᵢ · √(1−dᵢ²/rᵢ²)\n"
     "Realism pipeline: Poisson noise → PSF → gamma → misregistration jitter."),
    ("Step 5", "Output",
     "Fluoroscopy frame (H×W float32, GPU tensor) → RL policy input.\n"
     "Catheter state vector (tip xyz, orientation, insertion depth) → reward.\n"
     "Both outputs GPU-resident.  No host round-trip."),
]
colors = [NV_GREEN, ACCENT, RGBColor(0xE6,0x7E,0x22), RGBColor(0x9B,0x59,0xB6), NV_MID]
y = 1.98
for (num, title, body), col in zip(steps, colors):
    add_rect(sl, 0.25, y, 0.65, 0.95, fill=col)
    add_text(sl, num, 0.28, y+0.05, 0.6, 0.35, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 0.95, y+0.02, 2.8, 0.35, size=13, bold=True, color=col)
    add_text(sl, body,  0.95, y+0.35, 12.1, 0.65, size=10.5, color=NV_DARK)
    y += 1.05

add_text(sl,
    "Why this is fast:  CT volume never moves.  Only the N×8 catheter buffer changes per frame.  "
    "CUDA graph eliminates physics launch overhead.  ~63 FPS @ 256² / ~25 FPS @ 512² (A6000).",
    0.3, 7.1, 12.7, 0.3, size=11, bold=True, color=NV_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NEWTON UPSTREAM ALGORITHMS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Newton Upstream Algorithms — Ported to Isaac Lab")
add_text(sl,
    "XPBDRodSolver is a faithful port of Newton's SolverXPBDRod (PR #1981).  "
    "Feature parity maintained — research advances in Newton propagate directly into the Isaac Lab training environment.",
    0.3, 1.6, 12.7, 0.45, size=12, color=NV_MID, italic=True)

# Left column — algorithms
add_text(sl, "Core Algorithms Upstreamed", 0.3, 2.15, 6.2, 0.4, size=16, bold=True, color=NV_GREEN)

algo_rows = [
    ("Block-Thomas Direct Solve",
     "Tridiagonal block factorization for the XPBD constraint system.\n"
     "Linear in rod length O(N).  Numerically stable for stiff Cosserat rods.\n"
     "Runs entirely on GPU — no LAPACK, no CPU fallback."),
    ("Split-Thomas Solve",
     "Decoupled positional and rotational Thomas passes.\n"
     "Alternative backend selectable via XPBD_DIRECT_SOLVE_SPLIT_THOMAS."),
    ("Block-Jacobi Solve",
     "Parallel block-diagonal approximation.  Lower memory bandwidth.\n"
     "Useful when block-Thomas bandwidth is the bottleneck at very high N."),
    ("Banded Cholesky Solve",
     "Half-bandwidth KD=11, leading dimension LDAB=34.\n"
     "Exact Cholesky for small dense systems (N_dof ≤ 64)."),
    ("Tiled Cholesky (GPU tile)",
     "wp.tile_load for small dense systems — fits in shared memory tile.\n"
     "Avoids global memory round-trip for <64-DOF sub-problems."),
    ("Floor Collision",
     "Configurable restitution coefficient (0–1).\n"
     "Applied as a position correction constraint — parity with Newton upstream."),
]
y = 2.58
for title, body in algo_rows:
    add_rect(sl, 0.3, y, 0.08, 0.62, fill=NV_GREEN)
    add_text(sl, title, 0.45, y+0.01, 5.9, 0.28, size=12, bold=True, color=NV_DARK)
    add_text(sl, body, 0.45, y+0.28, 5.9, 0.38, size=10.5, color=NV_MID)
    y += 0.68

# Right column — vessel collision (XCath additions)
add_text(sl, "Vessel-Mesh Collision (XCathRodSolver)", 6.8, 2.15, 6.2, 0.4,
         size=16, bold=True, color=ACCENT)

vessel_items = [
    ("SDF Containment  (wp.mesh_query_point_sign_normal)",
     "BVH query per particle.  Signed distance to nearest surface point.\n"
     "Particles outside vessel lumen projected back along surface normal.\n"
     "sign_scale=1.0, target_phi=−0.001 m (1 mm inside wall)."),
    ("AABB Broadphase  (wp.mesh_query_aabb)",
     "Rod-segment bounding box queried against mesh triangle AABB tree.\n"
     "Candidate triangle faces → vertex-triangle + rod-segment-to-edge contacts.\n"
     "collision_iterations=2 per substep."),
    ("Track-Guided Insertion",
     "Non-tip particles projected to linear insertion axis.\n"
     "track_start + track_dir + track_length define the sheath entry axis.\n"
     "tip_num_edges controls how many distal segments are free to deflect."),
    ("Constraint Hooks",
     "_pre_constraints_hook: vessel SDF applied before XPBD stretch/bend.\n"
     "_post_constraints_hook: AABB edge contacts applied after XPBD solve.\n"
     "Extensible — additional constraint classes can be injected here."),
    ("CUDA-Graph Compatible",
     "All collision kernels are Warp kernels captured inside the same\n"
     "CUDA graph as the XPBD substep loop.  Zero CPU sync per substep."),
]
y = 2.58
for title, body in vessel_items:
    add_rect(sl, 6.8, y, 0.08, 0.78, fill=ACCENT)
    add_text(sl, title, 6.95, y+0.01, 6.1, 0.28, size=12, bold=True, color=NV_DARK)
    add_text(sl, body, 6.95, y+0.28, 6.1, 0.52, size=10.5, color=NV_MID)
    y += 0.84

add_rect(sl, 6.7, 2.15, 0.04, 5.2, fill=RGBColor(0xDD,0xDD,0xDD))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MULTI-ENVIRONMENT BATCHED PHYSICS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Multi-Environment Batched Physics — 512 Parallel Rods")
add_text(sl,
    "All N environments dispatched in a single Warp kernel launch and a single CUDA graph replay.  "
    "No per-environment CPU synchronization.",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

add_text(sl, "How Batching Works", 0.3, 2.1, 8.0, 0.38, size=16, bold=True, color=NV_GREEN)

batch_points = [
    "_BatchedWorkspace allocates flat contiguous GPU buffers:  "
    "positions[total_particles × 3],  orientations[total_particles × 4].",
    "rod_offsets[r] and edge_offsets[r] give each environment's slice — "
    "one kernel launch processes all rods in parallel without branching.",
    "CUDA-graph capture covers the entire multi-environment substep loop.  "
    "wp.ScopedCapture on first step() call; wp.capture_launch on all subsequent calls.",
    "GPU-side proximal control:  apply_proximal_control_gpu and set_root_orientation "
    "are Warp kernels — safe to capture inside the CUDA graph, zero CPU sync on hot path.",
    "Memory cost scales linearly with N (particle count × envs).  "
    "Compute cost is dominated by the block-Thomas solve — nearly O(N×envs) with the flat buffer layout.",
]
y = 2.58
for pt in batch_points:
    add_rect(sl, 0.3, y+0.08, 0.08, 0.22, fill=NV_GREEN)
    add_text(sl, pt, 0.45, y, 7.8, 0.42, size=12, color=NV_DARK)
    y += 0.50

add_text(sl, "Batched Rendering", 0.3, 5.25, 8.0, 0.38, size=16, bold=True, color=ACCENT)
render_points = [
    "renderDRR_forward_batched: single Slang dispatch.  dispatchThreadID.z indexes the environment.",
    "Per-environment catheter geometry packed into a flat StructuredBuffer<CatheterSegment> "
    "with per-env offset + count indices — different catheter per env, single kernel.",
    "L2 cache bottleneck at N>8: shared 3D μ-volume exceeds L2 capacity.  "
    "Sprint 2 fix: Texture2DArray caches one depth slice per env in L2.",
]
y = 5.68
for pt in render_points:
    add_rect(sl, 0.3, y+0.08, 0.08, 0.22, fill=ACCENT)
    add_text(sl, pt, 0.45, y, 7.8, 0.42, size=12, color=NV_DARK)
    y += 0.50

# Right — performance table
add_text(sl, "Performance  (NVIDIA A6000)", 8.6, 2.1, 4.5, 0.38,
         size=14, bold=True, color=NV_GREEN)
perf = [
    ("Physics — 1 env",     "1,300 Hz",  "✓ Target met"),
    ("Physics — 512 envs",  "~60 Hz",    "✓ Target met"),
    ("Render — 256²  1 env","~263 FPS",  "Interactive demo"),
    ("Render — 512²  1 env","~25 FPS",   "Standalone bench"),
    ("Render — 512²  N≤4",  "~25 FPS",   "Production today"),
    ("Render — 512²  N=8+", "Degrades",  "Sprint 2 fix"),
    ("Full loop (256²)",    "~63 FPS",   "Physics+render"),
]
table_slide(sl, ["Component", "Throughput", "Status"],
            perf, 8.6, 2.6, 4.6,
            col_widths=[1.9, 1.15, 1.55])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FLUOROSCOPY RENDERING — COMPOSITING PATHS (CORRECTED)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Fluoroscopy Rendering — Compositing Paths")

paths = [
    ("Slang GPU\n(unified loop)", NV_GREEN,
     "Fused volume + catheter Beer-Lambert in single GPU ray march.\n"
     "Catheter composited inline: μ_total(s) = μ_CT(s) + Σᵢ μᵢ·√(1−dᵢ²/rᵢ²)\n"
     "Static CT μ-volume on GPU texture.  Only catheter segment buffer changes per frame.\n"
     "Post-processing on same GPU timeline: Poisson → PSF → gamma → jitter.",
     "~63 FPS @ 256×256\n~25 FPS @ 512×512\n(NVIDIA A6000)"),
    ("CPU Beer-Lambert\n(reference path)", NV_MID,
     "NumPy ray-march for offline ground-truth generation and renderer validation.\n"
     "Per-segment cylinder chord:  t(d) = 2·√(r²−d²)  for d<r,  else 0.\n"
     "Final: I = I₀·exp(−Σᵢ μᵢ·chord_norm(u,v))\n"
     "Post-processing: scatter convolution → Poisson → PSF → gamma  (realism.py, CPU NumPy).\n"
     "NOT on the RL training path — for validation and visual QA only.",
     "~2–5 FPS\n(NumPy, CPU)\nReference only"),
    ("Isaac Lab USD quad\n(planned — Sprint 2)", RGBColor(0xAA,0xAA,0xAA),
     "DRR backdrop texture applied to a 3D USD quad in Omniverse viewport.\n"
     "Capsule prims for catheter markers.  RTX path tracing for photoreal scatter.\n"
     "Intended for visual QA and demo purposes — not for RL training data.",
     "Real-time (RTX)\nPlanned\nSprint 2"),
]

y = 2.05
for (name, col, desc, perf) in paths:
    add_rect(sl, 0.25, y, 2.2, 1.45, fill=col)
    add_text(sl, name, 0.3, y+0.1, 2.1, 0.7, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, perf, 0.3, y+0.82, 2.1, 0.58, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.5, y, 10.55, 1.45, fill=NV_GREY, line=col)
    add_text(sl, desc, 2.6, y+0.05, 10.35, 1.35, size=11.5, color=NV_DARK)
    y += 1.58

add_rect(sl, 0.25, 6.88, 12.8, 0.48, fill=NV_DARK)
add_text(sl,
    "Key property:  Catheter attenuation is additive in the exponent — "
    "μ_total = μ_CT + μ_catheter — so the catheter attenuates rather than occludes.  "
    "Self-crossings produce physically correct increased opacity.  "
    "Background anatomy remains visible through the catheter wire (exp(−μ) < 1, > 0).",
    0.35, 6.9, 12.6, 0.44, size=11, bold=True, color=NV_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BEER-LAMBERT CATHETER COMPOSITING — ALGORITHM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Beer-Lambert Catheter Compositing — Algorithm Detail")
add_text(sl,
    "Per-sample-point inline compositing inside the GPU ray march.  "
    "Same Beer-Lambert integral for CT and catheter — no separate compositing pass.",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

steps_bl = [
    ("1", "DRR Background",
     "Ray march through 3D CT μ-volume (GPU texture, static per patient).\n"
     "Fixed-step integration: I_DRR = I₀·exp(−∫μ_CT(s) ds)\n"
     "HU→μ piecewise linear map.  Dense bone ~0.3–0.5 mm⁻¹.  Air ~0.0 mm⁻¹."),
    ("2", "Catheter Segment Buffer",
     "Physics particle positions (N×3 float32) written to GPU segment buffer.\n"
     "Each segment: (xyz_proximal, xyz_distal, radius, μ) — 8 floats.\n"
     "Per-segment μ profile: tungsten marker 3.0, NiTi shaft 0.8, polymer tip 0.15, Pt marker 5.0."),
    ("3", "Inline Perpendicular Test",
     "At each ray sample point s, for each catheter segment i:\n"
     "  d_i = perpendicular distance from s to segment axis\n"
     "  if d_i < r_i:  contribute μᵢ·√(1−dᵢ²/rᵢ²) to running integral\n"
     "The √(1−d²/r²) factor is a chord-length weight — exact for a circular cross-section."),
    ("4", "Fused Integral",
     "μ_total(s) = μ_CT(s) + Σᵢ μᵢ·√(1−dᵢ²/rᵢ²)\n"
     "I_final = I₀·exp(−∫μ_total(s) ds)\n"
     "Catheter and CT share one exponent — no separate compositing pass, no alpha blending."),
    ("5", "Detector Realism",
     "Poisson quantum noise (photon count N_ph → Poisson(N_ph·I))\n"
     "PSF convolution (Gaussian σ=0.7 px — scintillator light spread)\n"
     "Gamma correction + misregistration jitter (σ_rot=0.05°, σ_trans=0.1 mm)"),
]

y = 2.05
for (num, title, body) in steps_bl:
    c = [NV_GREEN, ACCENT, RGBColor(0xE6,0x7E,0x22),
         RGBColor(0x9B,0x59,0xB6), RGBColor(0x16,0xA0,0x85)][int(num)-1]
    add_rect(sl, 0.25, y, 0.55, 0.95, fill=c)
    add_text(sl, num, 0.25, y+0.22, 0.55, 0.5, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 0.88, y+0.02, 3.0, 0.36, size=13, bold=True, color=c)
    add_text(sl, body, 0.88, y+0.38, 12.2, 0.6, size=10.5, color=NV_DARK)
    y += 1.02


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CT INGESTION PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Stage 1 — CT Ingestion Pipeline")
add_text(sl,
    "Runs once per patient anatomy.  Produces two GPU-resident assets that every downstream stage depends on.",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

pipeline_steps = [
    ("CT / CTA Input\n(DICOM / NIfTI)", "Raw 3D Hounsfield Unit volume from clinical scanner."),
    ("HU → μ Conversion", "Piecewise-linear map: HU[−1000,3000] → μ[0, 0.5] mm⁻¹.\nDense bone ~0.3–0.5.  Soft tissue ~0.02–0.05.  Air ~0.0."),
    ("μ Volume (3D tensor)", "Uploaded to GPU texture once.  Static for entire training session.\n→ Fed to Slang DRR renderer every frame."),
    ("Vessel Segmentation", "HU threshold on contrast-enhanced CT (vessels ~100–300 HU above baseline).\n→ Binary vessel mask."),
    ("Mesh Generation", "Marching Cubes / VTK on binary mask → polygon mesh.\nConverted to warp.Mesh (metres).  → Collision geometry for XCathRodSolver."),
    ("Centerline Extraction", "VMTK centerline skeleton from vessel mask.\n→ CenterlineGraph (nodes + edges).\n→ Dijkstra arrival map for temporal bolus dynamics."),
]

xs = [0.3, 2.8, 5.3, 7.8, 5.3, 7.8]
ys = [2.1, 2.1, 2.1, 2.1, 4.2, 4.2]
conn = [(0,1),(1,2),(2,3),(1,4),(4,5)]

for i, (title, body) in enumerate(pipeline_steps):
    col = NV_GREEN if i in (0,2,4) else ACCENT if i in (3,5) else NV_MID
    add_rect(sl, xs[i], ys[i], 2.3, 1.85, fill=NV_GREY, line=col)
    add_rect(sl, xs[i], ys[i], 2.3, 0.45, fill=col)
    add_text(sl, title, xs[i]+0.08, ys[i]+0.05, 2.15, 0.38, size=11, bold=True, color=WHITE)
    add_text(sl, body,  xs[i]+0.08, ys[i]+0.5,  2.15, 1.3,  size=10.5, color=NV_DARK)

add_text(sl,
    "Key:  The CT volume is a static 3D texture on the GPU.  "
    "Only the catheter geometry buffer changes frame-to-frame.  "
    "This is why the renderer can run at 25–63 FPS — no per-frame CT reconstruction.",
    0.3, 6.6, 12.7, 0.72, size=12, bold=True, color=NV_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DSA PIPELINE + BOLUS DYNAMICS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "DSA Pipeline + Temporal Bolus Dynamics")
add_text(sl,
    "4-step Digital Subtraction Angiography.  Gamma-variate contrast kinetics + Dijkstra arrival map.",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

dsa_steps = [
    ("Mask DRR", "Render anatomy without contrast (before injection).\nC-arm pose recorded precisely."),
    ("Contrast DRR", "Inject bolus:  Δμ(v,t) = μ_iodine · C(t − T(v)) · vessel_mask(v)\n"
     "C(t): gamma-variate  c_peak·(t/t_peak)^α·exp(α(1−t/t_peak))\n"
     "T(v): per-voxel arrival time from Dijkstra on CenterlineGraph."),
    ("Scatter + Jitter", "apply_scatter() adds Compton scatter haze.\n"
     "apply_misregistration(): σ_rot=0.05°, σ_trans=0.1 mm — simulates patient motion.\n"
     "Critical: introduces bone-edge residuals that monoenergetic rendering suppresses."),
    ("Subtract + Post", "DSA = contrast_DRR − mask_DRR  (in log domain).\n"
     "Post-process: normalize → gamma → clip → vessel boost (μ×factor on mask)."),
]
x = 0.25
for i, (title, body) in enumerate(dsa_steps):
    col = [NV_GREEN, ACCENT, RGBColor(0xE6,0x7E,0x22), RGBColor(0x9B,0x59,0xB6)][i]
    add_rect(sl, x, 2.1, 3.0, 2.2, fill=NV_GREY, line=col)
    add_rect(sl, x, 2.1, 3.0, 0.45, fill=col)
    add_text(sl, f"{i+1}. {title}", x+0.1, 2.13, 2.8, 0.38, size=13, bold=True, color=WHITE)
    add_text(sl, body, x+0.1, 2.6, 2.82, 1.65, size=10.5, color=NV_DARK)
    if i < 3:
        add_text(sl, "▶", x+3.04, 2.9, 0.25, 0.4, size=16, bold=True, color=col)
    x += 3.28

add_text(sl, "XCath Bolus Validation Results  (not yet merged into fluorosim)", 0.3, 4.48, 12.7, 0.38,
         size=15, bold=True, color=VALID_COL)
add_text(sl,
    "XCath validated Hagen-Poiseuille (power=2.0) and Feeding Trunk Excision (F-8) against real clinical DSA.  "
    "These results are algorithmically proven but the code changes are not yet merged into the fluorosim repo.  "
    "Current fluorosim uses Dijkstra with simple edge weights and multi-root injection.  Merge target: Sprint 2.",
    0.3, 4.9, 12.7, 0.65, size=11, color=NV_DARK, italic=True)

table_slide(sl,
    ["Metric", "Current fluorosim", "XCath validated (not merged)", "Real DSA"],
    [["Bolus delay",       "Not calibrated",  "1.90 s  (Hagen-Poiseuille)",  "2.13 s"],
     ["Peak delay error",  "—",               "11%  (was 2.6×)",             "—"],
     ["Distal MCA vel.",   "Not modelled",    "25% ICA  (viscous resist.)",  "Realistic"],
     ["Injection coverage","All roots",       "43.1%  (F-8 selective)",      "ICA territory"]],
    0.3, 5.63, 12.73,
    col_widths=[2.2, 2.55, 3.25, 2.73])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PERFORMANCE BASELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Performance Baseline  (NVIDIA A6000)")
add_text(sl,
    "Physics is solved.  Rendering is the diagnosed remaining constraint with a clear engineering path.",
    0.3, 1.6, 12.7, 0.35, size=13, color=NV_MID, italic=True)

table_slide(sl,
    ["Component", "Resolution / Scale", "Throughput", "Target", "Status"],
    [
        ["Single-env physics (XPBD)",  "20-seg rod",      "1,300 Hz",   ">1,000 Hz",  "✓ Met"],
        ["Multi-env physics",          "512 envs",        "~60 Hz",     ">60 Hz",     "✓ Met"],
        ["GPU DRR + catheter  (fused)","256×256  1 env",  "~263 FPS",   "<5 ms",      "✓ Met"],
        ["GPU DRR + catheter  (fused)","512×512  1 env",  "~25 FPS",    "<5 ms",      "In Progress"],
        ["GPU DRR batched",            "512×512  N≤4",    "~25 FPS",    ">60 Hz",     "In Progress"],
        ["GPU DRR batched",            "512×512  N>8",    "Degrades",   ">60 Hz",     "Sprint 2"],
        ["CPU Beer-Lambert",           "512×512  ref.",   "2–5 FPS",    "N/A",        "Reference only"],
        ["Full sim loop (interact.)",  "256×256  1 env",  "~63 FPS",    "30+ FPS",    "✓ Met"],
    ],
    0.25, 2.05, 12.83,
    col_widths=[3.2, 2.4, 1.7, 1.6, 1.93])

add_text(sl,
    "Sprint 2 bottleneck fix:  3D μ-volume → Texture2DArray — caches one depth slice per env in L2.  "
    "Eliminates global memory bandwidth saturation at N>8.  Projected throughput: >60 Hz @ 512 envs.",
    0.3, 6.3, 12.7, 0.55, size=12, bold=True, color=NV_GREEN)

add_text(sl,
    "Physics subsystem exceeds training throughput target with margin.  "
    "End-to-end loop is rendering-bound only.  No unknown unknowns on this slide.",
    0.3, 6.9, 12.7, 0.45, size=11, color=NV_MID, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RL TRAINING PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "RL Training Pipeline — State-Based (Sprint 1)")
add_text(sl,
    "PPO via RSL-RL.  512 parallel environments.  State observations today → pixel observations Sprint 2.",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

rl_rows = [
    ("Environment",
     "Multi-env production XPBD solver with gymnasium-compatible step/reset interface.\n"
     "Two control actions: advance (proximal push) + rotate — matching real procedure DOF.\n"
     "Reward: reduction in tip-to-target Euclidean distance.  Dense, well-conditioned."),
    ("RSL-RL Wrapper",
     "Standard VecEnv adapter.  RL algorithm fully decoupled from physics/rendering.\n"
     "Swapping to pixel obs or a different policy algorithm requires only obs tensor shape change."),
    ("PPO Config",
     "Hyperparameters tuned for catheter navigation: higher entropy (exploration at bifurcations),\n"
     "shorter GAE horizon (reduce variance from contact sequences)."),
    ("Training Scale",
     "512 parallel environments.  Policy sees 512 independent rollout trajectories per update.\n"
     "Physics runs at ~1,300 Hz/env inside a single CUDA graph — RL training is not physics-bound.\n"
     "Currently rendering-bound only if pixel obs are enabled."),
    ("Smoke Test",
     "Dependency-free validation path: reset → step → reward → done, no RSL-RL stack required.\n"
     "CI validation + headless server debugging without full training stack."),
]
y = 2.1
for i, (title, body) in enumerate(rl_rows):
    col = [NV_GREEN, ACCENT, RGBColor(0xE6,0x7E,0x22),
           RGBColor(0x9B,0x59,0xB6), NV_MID][i]
    add_rect(sl, 0.25, y, 2.2, 1.0, fill=col)
    add_text(sl, title, 0.28, y+0.25, 2.15, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.55, y, 10.5, 1.0, fill=NV_GREY, line=col)
    add_text(sl, body, 2.65, y+0.05, 10.3, 0.9, size=11, color=NV_DARK)
    y += 1.05

add_rect(sl, 0.25, 7.15, 12.8, 0.25, fill=RGBColor(0x1A,0x1A,0x1A))
add_text(sl,
    "State-based observations are scaffolding, not the target.  "
    "Sprint 2: pixel obs from Slang renderer → policy navigates using only the fluoroscopy frame — "
    "closing the loop from simulation to deployable clinical AI.",
    0.35, 7.17, 12.6, 0.22, size=10.5, bold=True, color=NV_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — INTERACTIVE UI DEMO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Stage 1 & 2 Demo — Interactive Fluoroscopy Simulator")
add_text(sl,
    "Browser-based UI integrating XPBD physics, Beer-Lambert rendering, vessel-mesh collision, and DSA in real time.",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

add_text(sl, "Control Panel", 0.3, 2.05, 4.5, 0.38, size=15, bold=True, color=NV_GREEN)
ctrl = [
    "C-arm projection: AP / LAO-30 / LAO-45 / RAO-30 / Lateral  (instant re-render)",
    "Advance speed (mm/s) slider  →  proximal boundary condition in XPBD solver",
    "Advance / Retract: 3 physics substeps per click",
    "Rotate CW / CCW: proximal torque, propagates along Cosserat rod",
    "Idle step: physics advance with zero control (gravity deformation)",
    "Reset: restore to initial straight configuration via Warp buffer write + CUDA graph invalidation",
    "DSA: 3-dispatch composite (bg + fat-catheter + catheter)",
]
y = 2.48
for c in ctrl:
    add_text(sl, f"• {c}", 0.35, y, 4.35, 0.3, size=11, color=NV_DARK)
    y += 0.3

add_text(sl, "Simulation Info (live)", 0.3, 4.72, 4.5, 0.35, size=15, bold=True, color=NV_GREEN)
info = [
    "Projection: current C-arm angle",
    "Tip (CT mm): catheter tip in CT volume coordinates",
    "Catheter bend (mm): max deviation from straight line — confirms vessel collision is active",
    "Physics step: substep time (ms)",
    "Render (GPU): DRR dispatch time (ms)",
    "Sim loop: total frame time + FPS",
]
y = 5.12
for c in info:
    add_text(sl, f"• {c}", 0.35, y, 4.35, 0.28, size=11, color=NV_DARK)
    y += 0.28

# Right — what the demo proves
add_text(sl, "Why Interactive Instead of PPO Rollout", 5.1, 2.05, 7.9, 0.38,
         size=15, bold=True, color=ACCENT)
table_slide(sl,
    ["Claim", "Why a pre-recorded rollout cannot prove it"],
    [
        ["Physics runs at interactive speed", "A recording can be played at any speed"],
        ["Vessel collision is active", "A scripted trajectory can be collision-free by design"],
        ["DRR re-renders at correct C-arm angle", "Policy rollouts typically use a fixed angle"],
        ["DSA + fluoro are spatially registered", "Hard to verify in a pre-rendered sequence"],
        ["Full pipeline runs in ~12 ms", "A video says nothing about throughput"],
    ],
    5.1, 2.5, 7.9,
    col_widths=[3.4, 4.5])

add_text(sl,
    "Launch:  conda activate isaaclab  →  python3 interactive_catheter_fluoro.py --ct-dir /tmp/patient_001 --share",
    0.3, 7.1, 12.7, 0.32, size=11, bold=True, color=NV_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — XCATH COLLABORATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "XCath Collaboration — Technical Outcomes")
add_text(sl,
    "5 algorithmic improvements validated against real clinical DSA data.  "
    "XCath embedded in the sprint — weekly technical sync with Jung-eun Park (AI/ML, XCATH Robotics).",
    0.3, 1.6, 12.7, 0.35, size=12, color=NV_MID, italic=True)

# Legend
for lx, lc, lbl in [(0.25, IMPL_COL,  "Implemented in fluorosim"),
                     (3.5,  VALID_COL, "Validated by XCath — not yet merged"),
                     (7.9,  ANALY_COL, "Analysis result — no code change")]:
    add_rect(sl, lx, 2.02, 0.22, 0.22, fill=lc)
    add_text(sl, lbl, lx+0.27, 1.99, 3.1, 0.28, size=11, color=NV_DARK)

# (title, header_col, status_col, status_label, body)
outcomes = [
    ("1. Hagen-Poiseuille\nBlood Flow Model",
     VALID_COL, VALID_COL, "Validated — not merged",
     "Murray's Law (power=0.5) → Hagen-Poiseuille (power=2.0)\n"
     "Bolus delay: 5.6s → 1.90s sim vs 2.13s real  (89% match, was 2.6× off)\n"
     "Distal MCA velocity: 71% ICA → 25% ICA  (viscous resistance)\n\n"
     "Not in fluorosim codebase — algorithm designed & validated\n"
     "by XCath on their data. Merge target: Sprint 2."),
    ("2. Selective Injection\n(Feeding Trunk Excision F-8)",
     VALID_COL, VALID_COL, "Validated — not merged",
     "Voronoi partition on centerline graph + excise non-selected trunks.\n"
     "Single-root coverage: 43.1% nodes — matches real ICA territory.\n"
     "Peak delay: 2.20s sim vs 2.13s real  (3% error)\n\n"
     "Not in fluorosim codebase — designed & validated by XCath.\n"
     "Merge target: Sprint 2."),
    ("3. Dispersion\nCorrection",
     VALID_COL, VALID_COL, "Validated — not merged",
     "β_eff = β₀ + k·arrival_time  (k=0.15) — FWHM widening.\n"
     "β_eff += α_res / max(r_voxel, 0.1) — thin vessels wider TDC.\n"
     "Distinguishability 0.607 (>0.3), asymmetry 2.55× (≥1.5×)\n\n"
     "Not in fluorosim codebase — designed & validated by XCath.\n"
     "Merge target: Sprint 2."),
    ("4. DeepDRR vs FluoroSim\nValidation",
     ANALY_COL, ANALY_COL, "Analysis only — no code",
     "SSIM=0.964, MAE=9.3%, PSNR=19.01 dB, Pearson=0.981\n"
     "Beam hardening cancels in DSA subtraction — monoenergetic\n"
     "rendering sufficient for DSA training data generation.\n\n"
     "External benchmarking by XCath — no code change needed.\n"
     "Finding informs Sprint 2 polyenergetic renderer priority."),
    ("5. Misregistration\nJitter",
     IMPL_COL, IMPL_COL, "Implemented ✓",
     "Controlled mask↔contrast shift: σ_rot=0.05°, σ_trans=0.1 mm.\n"
     "Introduces bone-edge residuals matching real DSA statistics.\n\n"
     "Implemented in fluorosim:\n"
     "  apply_misregistration()  — realism.py\n"
     "  misregistration_px param — DSASettings / dsa.py"),
]
x = 0.25
for (title, hdr_col, status_col, status_lbl, body) in outcomes:
    add_rect(sl, x, 2.3, 2.45, 4.5, fill=NV_GREY, line=hdr_col)
    add_rect(sl, x, 2.3, 2.45, 0.52, fill=hdr_col)
    add_text(sl, title, x+0.08, 2.32, 2.3, 0.46, size=11, bold=True, color=WHITE)
    # status badge
    add_rect(sl, x+0.08, 2.86, 2.3, 0.24, fill=status_col)
    add_text(sl, status_lbl, x+0.1, 2.87, 2.26, 0.22, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, body, x+0.08, 3.14, 2.32, 3.6, size=9.5, color=NV_DARK)
    x += 2.57

add_text(sl,
    "Strategic takeaway:  XCath provided clinical validation that turned the simulator into a calibrated training data generator.  "
    "Items 1–3 are algorithmically proven — merging them into fluorosim is a 2–3 day Sprint 2 task.",
    0.3, 6.92, 12.7, 0.45, size=11, bold=True, color=NV_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — COMPLETED DELIVERABLES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Sprint 1 — Completed Deliverables")

col1 = [
    ("Rendering", NV_GREEN, [
        "Slang GPU Beer-Lambert DRR  (512×512, ~25 FPS)",
        "Inline catheter compositing  (depth-correct, fused)",
        "Per-segment μ + radius  (marker / shaft / tip profile)",
        "Cone-beam magnification  (SID/SDD depth scaling)",
        "Batched multi-env DRR  (N envs, single Slang dispatch)",
        "Differentiable rendering  (6-DOF pose gradients, autodiff)",
    ]),
    ("Detector Realism + DSA", ACCENT, [
        "Poisson quantum noise  (photon-count model)",
        "Detector PSF convolution  (scintillator light spread)",
        "Gamma correction + normalize",
        "Misregistration jitter  ✓ implemented  (realism.py, dsa.py)",
        "4-step DSA pipeline  (mask→contrast→subtract→post)",
        "Temporal bolus dynamics  (gamma-variate + Dijkstra arrival map)",
        "Volumetric instrument injection  (Warp atomic_max, ~2 ms)",
        "9 vendor C-arm geometry presets",
        "Hagen-Poiseuille / F-8 / Dispersion  ⚠ validated, not yet merged",
    ]),
]
col2 = [
    ("Physics Solver", RGBColor(0xE6,0x7E,0x22), [
        "Cosserat rod XPBD  (block-Thomas direct solve)",
        "4 solver backends  (block-Thomas, split-Thomas, Jacobi, Cholesky)",
        "Multi-env batched physics  (_BatchedWorkspace, flat buffers)",
        "CUDA-graph capture  (wp.ScopedCapture + wp.capture_launch)",
        "GPU-side proximal control  (push + rotate, CUDA-graph safe)",
        "Floor collision  (configurable restitution, Newton parity)",
        "Vessel-mesh SDF containment  (wp.mesh_query_point_sign_normal)",
        "AABB broadphase + edge contacts  (wp.mesh_query_aabb)",
        "Track-guided insertion  (linear insertion axis constraint)",
    ]),
    ("Integration", RGBColor(0x9B,0x59,0xB6), [
        "CT ingestion pipeline  (HU→μ + vessel mask + wp.Mesh)",
        "Unified sim loop  (physics→render→RL obs, no CPU sync)",
        "Interactive Gradio UI  (browser-based, --share public link)",
        "RL training pipeline  (RSL-RL PPO, 512 envs, smoke test)",
        "Live catheter bend indicator  (vessel collision confirmation)",
        "pip-installable CLI  (xcath-fluoro entry point)",
    ]),
]

y = 2.05
for (title, col, items) in col1:
    add_text(sl, title, 0.25, y, 6.2, 0.35, size=14, bold=True, color=col)
    y2 = y + 0.38
    for it in items:
        add_rect(sl, 0.28, y2+0.07, 0.15, 0.15, fill=col)
        add_text(sl, it, 0.5, y2, 5.95, 0.3, size=11, color=NV_DARK)
        y2 += 0.28
    y = y2 + 0.12

y = 2.05
for (title, col, items) in col2:
    add_text(sl, title, 6.7, y, 6.4, 0.35, size=14, bold=True, color=col)
    y2 = y + 0.38
    for it in items:
        add_rect(sl, 6.73, y2+0.07, 0.15, 0.15, fill=col)
        add_text(sl, it, 6.95, y2, 6.1, 0.3, size=11, color=NV_DARK)
        y2 += 0.28
    y = y2 + 0.12

add_rect(sl, 6.6, 2.05, 0.04, 5.2, fill=RGBColor(0xDD,0xDD,0xDD))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
header_bar(sl, "Roadmap — Sprint 1 → Sprint 2 → Sprint 3+")

cols = [
    ("TODAY\n(Sprint 1 — Operational)", NV_GREEN, [
        "GPU Beer-Lambert DRR  @25 FPS / 512²",
        "XPBD physics  @1,300 Hz / env",
        "Vessel-mesh SDF + AABB collision",
        "DSA pipeline + temporal bolus",
        "9 C-arm vendor presets",
        "State-based PPO  (512 envs)",
        "Interactive demo  (Gradio UI)",
        "25 FPS end-to-end sim loop",
    ]),
    ("SPRINT 2\n(Renderer Scaling + Pixel Obs)", ACCENT, [
        "Texture2DArray renderer  (>60 Hz @ 512 envs)",
        "Pixel observations for RL  (raw fluoroscopy input)",
        "First end-to-end PPO policy  (fluoroscopy → action)",
        "Polyenergetic rendering  (beam hardening correction)",
        "GPU-side Poisson / PSF pipeline",
        "MAISI synthetic CT  (multi-patient training)",
        "Curriculum learning  (difficulty progression)",
    ]),
    ("SPRINT 3+\n(Clinical AI Pipeline)", RGBColor(0xE6,0x7E,0x22), [
        "2D/3D registration model  (end-to-end via autodiff)",
        "Vessel segmentation from fluoroscopy",
        "XCath robot integration  (policy → hardware commands)",
        "Holoscan IGX deployment  (<10 ms inference)",
        "Multi-patient validation  (5+ CTA datasets)",
        "Domain randomization study  (sim-to-real gap measurement)",
        "Clinical data paired evaluation",
    ]),
]
x = 0.25
for (title, col, items) in cols:
    add_rect(sl, x, 2.05, 4.18, 5.2, fill=NV_GREY, line=col)
    add_rect(sl, x, 2.05, 4.18, 0.65, fill=col)
    add_text(sl, title, x+0.1, 2.08, 3.98, 0.6, size=13, bold=True, color=WHITE)
    y = 2.78
    for it in items:
        add_rect(sl, x+0.1, y+0.09, 0.12, 0.12, fill=col)
        add_text(sl, it, x+0.28, y, 3.82, 0.28, size=11, color=NV_DARK)
        y += 0.3
    x += 4.4

add_text(sl,
    "This project is data infrastructure for robotic catheter AI.  "
    "Sprint 1 built the foundation.  Sprint 2 closes the sim-to-real loop.  "
    "Sprint 3 deploys to a clinical robot.",
    0.3, 7.12, 12.7, 0.28, size=12, bold=True, color=NV_GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"Saved: {OUT}")
